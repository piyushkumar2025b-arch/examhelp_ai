"""
presentation_engine.py — AI Presentation Builder (Powerup v2)
Per-slide internet data via DuckDuckGo, Unsplash images embedded in PPTX,
parallel enrichment via concurrent.futures, per-slide improve, PDF/MD export.
"""

import streamlit as st
import io
import re
import json
import requests
import concurrent.futures
import hashlib
from utils import ai_engine

THEMES = {
    "Professional": {
        "bg_color": "FFFFFF", "title_color": "1a1a2e", "text_color": "333333",
        "accent": "2563eb", "font": "Calibri",
    },
    "Academic": {
        "bg_color": "FAFAF5", "title_color": "1e3a5f", "text_color": "2c3e50",
        "accent": "1e3a5f", "font": "Times New Roman",
    },
    "Creative": {
        "bg_color": "FFF5E6", "title_color": "e74c3c", "text_color": "2c3e50",
        "accent": "e74c3c", "font": "Segoe UI",
    },
    "Dark": {
        "bg_color": "1a1a2e", "title_color": "a78bfa", "text_color": "e0e0ff",
        "accent": "7c6af7", "font": "Segoe UI",
    },
    "Minimal": {
        "bg_color": "FFFFFF", "title_color": "111111", "text_color": "444444",
        "accent": "000000", "font": "Helvetica",
    },
}

AUDIENCE_TYPES = [
    "General Audience", "Technical Experts", "Business Executives",
    "Students", "Academic Conference", "Investors/Pitch Deck",
]

TONES = ["Professional", "Casual & Engaging", "Academic & Formal",
         "Inspirational", "Data-Driven", "Storytelling"]


def _hex_to_rgb(hex_color: str):
    h = hex_color.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def generate_slide_structure(topic: str, num_slides: int, audience: str, tone: str, notes_text: str = "") -> list:
    context = f"\nAdditional context/notes:\n{notes_text[:2000]}" if notes_text else ""

    prompt = f"""Create a presentation outline for: "{topic}"
Number of content slides: {num_slides}
Audience: {audience}
Tone: {tone}
{context}

Return a JSON array of slide objects. Each slide:
{{
  "slide_number": 1,
  "type": "title|agenda|content|summary|qna",
  "title": "Slide Title",
  "bullets": ["Key point 1", "Key point 2", "Key point 3"],
  "speaker_notes": "What to say during this slide",
  "image_keyword": "keyword for finding a relevant image"
}}

Structure:
- Slide 1: Title slide (type: "title")
- Slide 2: Agenda/Overview (type: "agenda")
- Slides 3 to {num_slides}: Content slides (type: "content")
- Second to last: Summary/Conclusion (type: "summary")
- Last: Q&A slide (type: "qna")

Make bullets concise (max 15 words each). Include real statistics or data points where possible.
Return ONLY the JSON array."""

    try:
        result = ai_engine.generate(
            prompt=prompt,
            model="llama-3.3-70b-versatile",
            json_mode=True,
            max_tokens=4096,
            temperature=0.4,
        )
        data = json.loads(result) if isinstance(result, str) else result
        if isinstance(data, list):
            return data
    except Exception:
        pass

    basic = [
        {"slide_number": 1, "type": "title", "title": topic, "bullets": [audience, tone], "speaker_notes": f"Welcome to this presentation on {topic}.", "image_keyword": topic.split()[0]},
        {"slide_number": 2, "type": "agenda", "title": "Agenda", "bullets": ["Overview", "Key Points", "Conclusion"], "speaker_notes": "Here's what we'll cover.", "image_keyword": "agenda"},
    ]
    for i in range(3, num_slides + 1):
        basic.append({"slide_number": i, "type": "content", "title": f"Section {i-2}", "bullets": [f"Point {j}" for j in range(1, 4)], "speaker_notes": "", "image_keyword": topic.split()[0]})
    basic.append({"slide_number": num_slides + 1, "type": "summary", "title": "Summary", "bullets": ["Key takeaways"], "speaker_notes": "Let me summarize.", "image_keyword": "summary"})
    basic.append({"slide_number": num_slides + 2, "type": "qna", "title": "Questions?", "bullets": ["Thank you for your attention"], "speaker_notes": "Open floor for Q&A.", "image_keyword": "questions"})
    return basic


def fetch_slide_internet_data(slide_title: str, topic: str) -> str:
    cache_key = f"_pres_web_{__import__('hashlib').md5((slide_title+topic).encode()).hexdigest()[:10]}"
    cached = st.session_state.get(cache_key)
    if cached:
        return cached
    try:
        from duckduckgo_search import DDGS
        query = f"{topic} {slide_title} statistics facts"
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=3))
        if results:
            snippet = results[0].get("body", "")[:300]
            st.session_state[cache_key] = snippet
            return snippet
    except Exception:
        pass
    return ""


def parallel_enrich_all_slides(slides: list, topic: str) -> list:
    import concurrent.futures
    def _enrich_one(args):
        idx, slide = args
        enriched = enrich_slide(slide, topic)
        if slide.get("type") not in ("title", "qna"):
            web_fact = fetch_slide_internet_data(slide.get("title", ""), topic)
            if web_fact:
                enriched["speaker_notes"] = enriched.get("speaker_notes", "") + f"\n\n📰 Live: {web_fact[:200]}"
                enriched["_web_snippet"] = web_fact
        return idx, enriched

    result = list(slides)
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as ex:
        for idx, enriched in ex.map(_enrich_one, enumerate(slides)):
            result[idx] = enriched
    return result


def download_image_bytes(keyword: str) -> bytes:
    try:
        url = f"https://source.unsplash.com/800x600/?{keyword.replace(' ', '+')}"
        r = requests.get(url, timeout=8, allow_redirects=True)
        if r.status_code == 200 and r.content:
            return r.content
    except Exception:
        pass
    return b""


def enrich_slide(slide: dict, topic: str) -> dict:
    if slide.get("type") in ("title", "qna"):
        return slide

    prompt = f"""Enrich this presentation slide with real data.
Topic: {topic}
Slide title: {slide.get('title', '')}
Current bullets: {json.dumps(slide.get('bullets', []))}

Add real statistics, facts, or expert quotes to each bullet point.
Return a JSON object: {{"bullets": ["enriched point 1", "enriched point 2", ...], "speaker_notes": "enhanced notes"}}
Keep each bullet under 20 words. Include source references where possible."""

    try:
        result = ai_engine.generate(
            prompt=prompt,
            model="llama-3.3-70b-versatile",
            json_mode=True,
            max_tokens=1024,
            temperature=0.3,
        )
        data = json.loads(result) if isinstance(result, str) else result
        if isinstance(data, dict):
            if "bullets" in data:
                slide["bullets"] = data["bullets"]
            if "speaker_notes" in data:
                slide["speaker_notes"] = data["speaker_notes"]
    except Exception:
        pass
    return slide


def improve_slide(slide: dict, instruction: str) -> dict:
    prompt = f"""Improve this presentation slide based on the instruction.
Current title: {slide.get('title', '')}
Current bullets: {json.dumps(slide.get('bullets', []))}
Current notes: {slide.get('speaker_notes', '')}

Instruction: {instruction}

Return a JSON object: {{"title": "...", "bullets": [...], "speaker_notes": "..."}}"""

    try:
        result = ai_engine.generate(
            prompt=prompt,
            model="llama-3.3-70b-versatile",
            json_mode=True,
            max_tokens=1024,
            temperature=0.5,
        )
        data = json.loads(result) if isinstance(result, str) else result
        if isinstance(data, dict):
            slide.update(data)
    except Exception:
        pass
    return slide


def get_unsplash_url(keyword: str, width: int = 800, height: int = 600) -> str:
    return f"https://source.unsplash.com/{width}x{height}/?{keyword.replace(' ', '+')}"


def generate_pptx(slides: list, theme_name: str, footer_text: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    theme = THEMES.get(theme_name, THEMES["Professional"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_rgb = _hex_to_rgb(theme["bg_color"])
    title_rgb = _hex_to_rgb(theme["title_color"])
    text_rgb = _hex_to_rgb(theme["text_color"])
    accent_rgb = _hex_to_rgb(theme["accent"])

    for slide_data in slides:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("speaker_notes", "")

        if slide_type == "title":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_rgb)
            p.font.name = theme["font"]
            p.alignment = PP_ALIGN.CENTER

            if bullets:
                sub = tf.add_paragraph()
                sub.text = " · ".join(bullets[:3])
                sub.font.size = Pt(18)
                sub.font.color.rgb = RGBColor(*text_rgb)
                sub.font.name = theme["font"]
                sub.alignment = PP_ALIGN.CENTER
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_rgb)
            p.font.name = theme["font"]

            accent_line = slide.shapes.add_shape(
                1, Inches(0.8), Inches(1.4), Inches(2), Inches(0.04)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = RGBColor(*accent_rgb)
            accent_line.line.fill.background()

            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11), Inches(5))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            for j, bullet in enumerate(bullets):
                p = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(*text_rgb)
                p.font.name = theme["font"]
                p.space_after = Pt(8)

        if footer_text:
            ft = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.4))
            fp = ft.text_frame.paragraphs[0]
            fp.text = footer_text
            fp.font.size = Pt(9)
            fp.font.color.rgb = RGBColor(150, 150, 150)
            fp.font.name = theme["font"]
            fp.alignment = PP_ALIGN.RIGHT

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def generate_pptx_with_images(slides: list, theme_name: str, footer_text: str = "", embed_images: bool = True) -> bytes:
    """Generate PPTX with Unsplash images downloaded and embedded per slide."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io as _io

    theme = THEMES.get(theme_name, THEMES["Professional"])
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_rgb = _hex_to_rgb(theme["bg_color"])
    title_rgb = _hex_to_rgb(theme["title_color"])
    text_rgb = _hex_to_rgb(theme["text_color"])
    accent_rgb = _hex_to_rgb(theme["accent"])

    for slide_data in slides:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_rgb)

        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("speaker_notes", "")
        img_keyword = slide_data.get("image_keyword", "")

        img_embedded = False
        content_width = Inches(11)

        if embed_images and img_keyword and slide_type in ("content", "agenda", "summary"):
            img_bytes = download_image_bytes(img_keyword)
            if img_bytes:
                try:
                    img_stream = _io.BytesIO(img_bytes)
                    slide.shapes.add_picture(img_stream, Inches(9.5), Inches(1.5), Inches(3.5), Inches(2.6))
                    img_embedded = True
                    content_width = Inches(8.5)
                except Exception:
                    pass

        if slide_type == "title":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_rgb)
            p.font.name = theme["font"]
            p.alignment = PP_ALIGN.CENTER
            if bullets:
                sub = tf.add_paragraph()
                sub.text = " · ".join(bullets[:3])
                sub.font.size = Pt(18)
                sub.font.color.rgb = RGBColor(*text_rgb)
                sub.font.name = theme["font"]
                sub.alignment = PP_ALIGN.CENTER
        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), content_width, Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_rgb)
            p.font.name = theme["font"]

            accent_line = slide.shapes.add_shape(1, Inches(0.8), Inches(1.4), Inches(2), Inches(0.04))
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = RGBColor(*accent_rgb)
            accent_line.line.fill.background()

            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), content_width, Inches(5))
            ctf = content_box.text_frame
            ctf.word_wrap = True
            for j, bullet in enumerate(bullets):
                p2 = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
                p2.text = f"• {bullet}"
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(*text_rgb)
                p2.font.name = theme["font"]
                p2.space_after = Pt(8)

        if footer_text:
            ft = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.4))
            fp = ft.text_frame.paragraphs[0]
            fp.text = footer_text
            fp.font.size = Pt(9)
            fp.font.color.rgb = RGBColor(150, 150, 150)
            fp.font.name = theme["font"]
            fp.alignment = PP_ALIGN.RIGHT

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def generate_markdown_outline(slides: list) -> str:
    lines = ["# Presentation Outline\n"]
    for s in slides:
        lines.append(f"\n## Slide {s.get('slide_number', '')}: {s.get('title', '')}")
        lines.append(f"*Type: {s.get('type', 'content')}*\n")
        for b in s.get("bullets", []):
            lines.append(f"- {b}")
        if s.get("speaker_notes"):
            lines.append(f"\n> **Speaker Notes:** {s['speaker_notes']}")
    return "\n".join(lines)


def render_presentation_builder():
    st.markdown("""
<style>
.pres-header{background:linear-gradient(135deg,#0a1020 0%,#081018 100%);border:1px solid #2a4a6b;border-radius:16px;padding:28px 32px;margin-bottom:20px;}
.pres-title{font-size:2rem;font-weight:900;background:linear-gradient(135deg,#60a5fa,#a78bfa);-webkit-background-clip:text;-webkit-text-fill-color:transparent;background-clip:text;margin:0 0 4px;}
.pres-sub{font-size:.9rem;color:#9090b8;}
.slide-card{background:var(--bg2-glass);border:1px solid var(--bd-glass);border-radius:14px;padding:20px;margin:10px 0;transition:all .25s;backdrop-filter:blur(12px);}
.slide-card:hover{border-color:var(--accent-bd);transform:translateY(-2px);}
.slide-number{display:inline-flex;align-items:center;justify-content:center;width:28px;height:28px;background:var(--accent);color:white;border-radius:50%;font-size:.75rem;font-weight:700;}
.slide-type-badge{display:inline-flex;align-items:center;gap:4px;background:var(--accent-bg);border:1px solid var(--accent-bd);border-radius:8px;padding:2px 8px;font-size:.7rem;color:var(--accent);font-weight:600;text-transform:uppercase;}
</style>
""", unsafe_allow_html=True)

    st.markdown("""
<div class="pres-header">
  <div class="pres-title">🎯 AI Presentation Builder</div>
  <div class="pres-sub">Generate complete slide decks with real data · Export as .pptx · AI-enriched content · Unsplash images</div>
</div>
""", unsafe_allow_html=True)

    if "pres_slides" not in st.session_state:
        st.session_state.pres_slides = []

    tab_create, tab_edit, tab_export = st.tabs(["🛠️ Create", "✏️ Edit & Preview", "📥 Export"])

    with tab_create:
        topic = st.text_input("📌 Presentation Topic",
                              placeholder="e.g. The Future of AI in Healthcare",
                              key="pres_topic")

        c1, c2, c3 = st.columns(3)
        with c1:
            num_slides = st.slider("Number of content slides", 3, 20, 8, key="pres_num")
        with c2:
            audience = st.selectbox("👥 Audience", AUDIENCE_TYPES, key="pres_audience")
        with c3:
            tone = st.selectbox("🎭 Tone", TONES, key="pres_tone")

        c4, c5 = st.columns(2)
        with c4:
            theme = st.selectbox("🎨 Theme", list(THEMES.keys()), key="pres_theme")
        with c5:
            footer = st.text_input("Footer text (optional)", placeholder="Company/University name",
                                   key="pres_footer")

        notes_file = st.file_uploader("📄 Upload notes/PDF (optional)", type=["txt", "pdf", "md"],
                                      key="pres_notes_upload")
        notes_text = ""
        if notes_file:
            try:
                if notes_file.name.endswith('.pdf'):
                    from utils.pdf_handler import extract_text_from_pdf
                    notes_text = extract_text_from_pdf(notes_file)
                else:
                    notes_text = notes_file.read().decode('utf-8', errors='replace')
            except Exception:
                notes_text = ""

        enrich = st.checkbox("🔍 Enrich slides with real internet data", value=True, key="pres_enrich")

        b1, b2 = st.columns([3, 1])
        with b1:
            gen_btn = st.button("🚀 Generate Presentation", type="primary", use_container_width=True,
                                disabled=not topic.strip(), key="pres_gen")
        with b2:
            if st.button("💬 Back", use_container_width=True, key="pres_back"):
                st.session_state.app_mode = "chat"
                st.rerun()

        if gen_btn and topic.strip():
            with st.spinner("🎯 Generating slide structure..."):
                slides = generate_slide_structure(topic, num_slides, audience, tone, notes_text)
                st.session_state.pres_slides = slides

            if enrich:
                with st.spinner("⚡ Enriching all slides in parallel with live internet data..."):
                    st.session_state.pres_slides = parallel_enrich_all_slides(
                        st.session_state.pres_slides, topic
                    )

            st.success(f"✅ Generated {len(st.session_state.pres_slides)} slides!")
            st.rerun()

    with tab_edit:
        if not st.session_state.pres_slides:
            st.info("Generate a presentation first in the Create tab.")
        else:
            slides = st.session_state.pres_slides

            for i, slide in enumerate(slides):
                slide_type = slide.get("type", "content")
                type_icons = {"title": "🏷️", "agenda": "📋", "content": "📝", "summary": "✅", "qna": "❓"}
                icon = type_icons.get(slide_type, "📄")

                with st.expander(f"**Slide {i+1}** — {icon} {slide.get('title', 'Untitled')}", expanded=False):
                    new_title = st.text_input("Title", value=slide.get("title", ""), key=f"pres_title_{i}")
                    new_bullets = st.text_area(
                        "Bullets (one per line)",
                        value="\n".join(slide.get("bullets", [])),
                        height=120,
                        key=f"pres_bullets_{i}"
                    )
                    new_notes = st.text_area(
                        "Speaker Notes",
                        value=slide.get("speaker_notes", ""),
                        height=80,
                        key=f"pres_notes_{i}"
                    )

                    if new_title != slide.get("title", ""):
                        st.session_state.pres_slides[i]["title"] = new_title
                    if new_bullets != "\n".join(slide.get("bullets", [])):
                        st.session_state.pres_slides[i]["bullets"] = [b for b in new_bullets.split("\n") if b.strip()]
                    if new_notes != slide.get("speaker_notes", ""):
                        st.session_state.pres_slides[i]["speaker_notes"] = new_notes

                    img_keyword = slide.get("image_keyword", "")
                    if img_keyword:
                        img_url = get_unsplash_url(img_keyword, 400, 250)
                        try:
                            st.image(img_url, caption=f"Image: {img_keyword}", use_container_width=True)
                        except Exception:
                            pass

                    ic1, ic2 = st.columns(2)
                    with ic1:
                        improve_instruction = st.text_input("Improve instruction", key=f"pres_improve_{i}",
                                                            placeholder="e.g. Add more statistics")
                    with ic2:
                        if st.button("✨ Improve", key=f"pres_improve_btn_{i}", use_container_width=True):
                            if improve_instruction:
                                with st.spinner("Improving slide..."):
                                    st.session_state.pres_slides[i] = improve_slide(slide, improve_instruction)
                                st.rerun()

    with tab_export:
        if not st.session_state.pres_slides:
            st.info("Generate a presentation first.")
        else:
            slides = st.session_state.pres_slides
            theme_name = st.session_state.get("pres_theme", "Professional")
            footer_text = st.session_state.get("pres_footer", "")

            st.markdown(f"**{len(slides)} slides** · Theme: {theme_name}")

            e1, e2, e3 = st.columns(3)

            embed_imgs = st.checkbox("🖼️ Embed Unsplash images in slides", value=True, key="pres_embed_imgs")

            with e1:
                try:
                    if embed_imgs:
                        with st.spinner("Downloading slide images..."):
                            pptx_bytes = generate_pptx_with_images(slides, theme_name, footer_text, embed_images=True)
                    else:
                        pptx_bytes = generate_pptx(slides, theme_name, footer_text)
                    st.download_button(
                        "📥 Download .PPTX" + (" (with images)" if embed_imgs else ""),
                        pptx_bytes,
                        file_name="presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key="pres_dl_pptx",
                    )
                except Exception as e:
                    st.error(f"PPTX generation error: {e}")

            with e2:
                md_content = generate_markdown_outline(slides)
                st.download_button(
                    "📥 Download Markdown",
                    md_content,
                    file_name="presentation_outline.md",
                    mime="text/markdown",
                    use_container_width=True,
                    key="pres_dl_md",
                )

            with e3:
                try:
                    from fpdf import FPDF
                    pdf = FPDF()
                    pdf.set_auto_page_break(auto=True, margin=15)
                    for s in slides:
                        pdf.add_page()
                        pdf.set_font("Helvetica", "B", 24)
                        pdf.cell(0, 20, s.get("title", ""), ln=True, align="C")
                        pdf.set_font("Helvetica", "", 14)
                        for b in s.get("bullets", []):
                            pdf.cell(0, 10, f"  - {b}", ln=True)
                    pdf_bytes = pdf.output()
                    st.download_button(
                        "📥 Download PDF",
                        pdf_bytes,
                        file_name="presentation.pdf",
                        mime="application/pdf",
                        use_container_width=True,
                        key="pres_dl_pdf",
                    )
                except Exception as e:
                    st.error(f"PDF generation error: {e}")

            with st.expander("📋 Full Markdown Preview"):
                st.markdown(generate_markdown_outline(slides))
