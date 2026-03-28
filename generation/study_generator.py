"""study_generator.py — REPAIRED"""
import io
from fpdf import FPDF
from docx import Document
from pptx import Presentation
import markdown2

class StudyGenerator:
    @staticmethod
    def generate_pdf(title, content):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        safe_content = content.encode("latin-1", "ignore").decode("latin-1")
        pdf.multi_cell(0, 10, safe_content)
        return bytes(pdf.output(dest="S"))

    @staticmethod
    def generate_docx(title, content):
        doc = Document()
        doc.add_heading(title, 0)
        doc.add_paragraph(content)
        stream = io.BytesIO()
        doc.save(stream)
        return stream.getvalue()

    @staticmethod
    def generate_ppt(title, content):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        stream = io.BytesIO()
        prs.save(stream)
        return stream.getvalue()
