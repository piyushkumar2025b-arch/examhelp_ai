"""
smart_reader_engine.py
═══════════════════════════════════════════════════════════════
ExamHelp AI — Smart AI Reader 🔬
Supports: PDF · PNG/JPG/WEBP · MP4/Video · PPT/PPTX

AI Functions per format:
  PDF  → Summarize · Q&A · Flashcards · Quiz · Keywords · Translate
          · Chapter-by-chapter · Timeline · Mind-map text · Critique
          · Simplify · Explain like 5 · Action items · Reading analytics
  IMG  → Describe · OCR / Extract text · Explain diagram · Chart-to-data
          · Handwriting reader · Formula solver · Sentiment/mood · Object list
          · Alt-text generator · Accessibility audit · Artistic analysis
  PPT  → Slide-by-slide summary · Speaker notes · Quiz from slides
          · Executive summary · Audience Q&A · Slide critique · Storytelling arc
  VID  → Frame-by-frame description · Script extraction · Key moments
          · Action items · Transcript summary · Topic timeline
═══════════════════════════════════════════════════════════════
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
import time
import random
from typing import List, Tuple, Optional

import streamlit as st

# ── lazy-safe imports ────────────────────────────────────────
try:
    import fitz                       # PyMuPDF  – PDF
except ImportError:
    fitz = None

try:
    from PIL import Image as PILImage  # images
except ImportError:
    PILImage = None

try:
    from pptx import Presentation     # PPT/PPTX
except ImportError:
    Presentation = None

try:
    import cv2                         # video frames
    import numpy as np
except ImportError:
    cv2 = None
    np = None

# ── AI backend ───────────────────────────────────────────────
from utils.ai_engine import generate, vision_generate


# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════

def _pil_to_b64(img) -> Tuple[str, str]:
    """Convert PIL Image → (base64-string, mime-type)."""
    buf = io.BytesIO()
    fmt = img.format if img.format else "JPEG"
    img.save(buf, format=fmt)
    b64 = base64.b64encode(buf.getvalue()).decode()
    mime = f"image/{fmt.lower()}"
    return b64, mime


def _bytes_to_b64(data: bytes, mime: str = "image/jpeg") -> str:
    return base64.b64encode(data).decode()


def _call_ai(prompt: str, system: str = "", max_tokens: int = 2048,
             temperature: float = 0.7) -> str:
    try:
        return generate(
            prompt=prompt,
            system=system,
            temperature=temperature,
            max_tokens=max_tokens,
        )
    except Exception as e:
        return f"⚠️ AI error: {e}"


def _call_vision(prompt: str, b64: str, mime: str,
                 max_tokens: int = 2048) -> str:
    try:
        return vision_generate(
            prompt=prompt,
            image_b64=b64,
            mime=mime,
            max_tokens=max_tokens,
        )
    except Exception as e:
        return f"⚠️ Vision AI error: {e}"


def _extract_pdf_text(data: bytes) -> Tuple[str, dict, List[str]]:
    """Return (full_text, metadata, page_texts[])."""
    if not fitz:
        return "PyMuPDF not installed.", {}, []
    try:
        doc = fitz.open(stream=data, filetype="pdf")
        meta = doc.metadata or {}
        meta["page_count"] = doc.page_count
        page_texts = []
        total_words = 0
        for page in doc:
            t = page.get_text().strip()
            page_texts.append(t)
            total_words += len(t.split())
        meta["word_count"] = total_words
        meta["reading_min"] = max(1, round(total_words / 250))
        full = "\n\n".join(
            f"[Page {i+1}]\n{t}" for i, t in enumerate(page_texts) if t
        )
        return full[:80_000], meta, page_texts
    except Exception as e:
        return f"Error: {e}", {}, []


def _extract_pptx(data: bytes) -> Tuple[List[dict], str]:
    """Return (slides[], full_text). slide = {num, title, body, notes}."""
    if not Presentation:
        return [], "python-pptx not installed."
    try:
        prs = Presentation(io.BytesIO(data))
        slides = []
        all_text = []
        for i, slide in enumerate(prs.slides):
            title = ""
            body_parts = []
            notes_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.shape_type == 13:   # picture
                        continue
                    if shape.name.lower().startswith("title") or (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                        and shape.placeholder_format.idx == 0
                    ):
                        title = shape.text.strip()
                    else:
                        body_parts.append(shape.text.strip())
            if slide.has_notes_slide:
                tf = slide.notes_slide.notes_text_frame
                notes_text = tf.text.strip() if tf else ""
            body = "\n".join(b for b in body_parts if b)
            slides.append({
                "num": i + 1,
                "title": title or f"Slide {i+1}",
                "body": body,
                "notes": notes_text,
            })
            all_text.append(f"[Slide {i+1}: {title}]\n{body}")
        return slides, "\n\n".join(all_text)
    except Exception as e:
        return [], f"Error reading PPT: {e}"


def _extract_video_frames(data: bytes, max_frames: int = 6) -> List[Tuple[str, str, float]]:
    """Return list of (b64, mime, timestamp_sec) for sampled frames."""
    if not cv2 or not np:
        return []
    try:
        tmp_path = "tmp_sr_video.mp4"
        with open(tmp_path, "wb") as f:
            f.write(data)
        cap = cv2.VideoCapture(tmp_path)
        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps   = cap.get(cv2.CAP_PROP_FPS) or 30
        dur   = total / fps
        frames = []
        indices = [int(i * total / max_frames) for i in range(max_frames)]
        for idx in indices:
            cap.set(cv2.CAP_PROP_POS_FRAMES, idx)
            ret, frame = cap.read()
            if not ret:
                continue
            ts = idx / fps
            _, buf = cv2.imencode(".jpg", frame)
            b64 = base64.b64encode(buf.tobytes()).decode()
            frames.append((b64, "image/jpeg", ts))
        cap.release()
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        return frames
    except Exception:
        return []


# ═══════════════════════════════════════════════════════════
# CSS STYLES
# ═══════════════════════════════════════════════════════════

_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=JetBrains+Mono:wght@400;700&family=Syne:wght@400;700;800&family=Plus+Jakarta+Sans:wght@400;500;600;700;800&display=swap');

/* ═══ GLOBAL ═══ */
.sr-root { font-family: 'Plus Jakarta Sans', 'Inter', sans-serif; }

/* ═══ ANIMATED BACKGROUND ORBS ═══ */
.sr-orb {
    position: fixed; border-radius: 50%; pointer-events: none;
    filter: blur(120px); z-index: 0;
}
.sr-orb-1 {
    width: 500px; height: 500px;
    background: radial-gradient(circle, rgba(99,102,241,0.18), transparent 70%);
    top: -150px; left: -100px;
    animation: srOrb1 14s ease-in-out infinite;
}
.sr-orb-2 {
    width: 400px; height: 400px;
    background: radial-gradient(circle, rgba(6,182,212,0.14), transparent 70%);
    bottom: 0; right: -100px;
    animation: srOrb2 18s ease-in-out infinite;
}
.sr-orb-3 {
    width: 300px; height: 300px;
    background: radial-gradient(circle, rgba(139,92,246,0.12), transparent 70%);
    top: 40%; left: 40%;
    animation: srOrb3 11s ease-in-out infinite;
}
@keyframes srOrb1 { 0%,100%{transform:translate(0,0);} 50%{transform:translate(60px,-40px);} }
@keyframes srOrb2 { 0%,100%{transform:translate(0,0);} 50%{transform:translate(-50px,30px);} }
@keyframes srOrb3 { 0%,100%{transform:translate(0,0) scale(1);} 50%{transform:translate(30px,-20px) scale(1.1);} }

/* ═══ HERO ═══ */
.sr-hero {
    position: relative; overflow: hidden;
    border-radius: 28px;
    padding: 48px 36px 42px;
    margin-bottom: 28px;
    text-align: center;
    background:
        radial-gradient(ellipse at 20% 0%, rgba(99,102,241,0.2) 0%, transparent 60%),
        radial-gradient(ellipse at 80% 100%, rgba(6,182,212,0.15) 0%, transparent 60%),
        rgba(10,14,30,0.85);
    border: 1px solid rgba(99,102,241,0.22);
    backdrop-filter: blur(24px);
    box-shadow:
        0 0 0 1px rgba(255,255,255,0.04) inset,
        0 40px 100px rgba(0,0,0,0.5),
        0 0 80px rgba(99,102,241,0.08);
    animation: srFadeUp 0.7s cubic-bezier(0.16,1,0.3,1) both;
}
@keyframes srFadeUp {
    from { opacity:0; transform: translateY(24px) scale(0.98); }
    to   { opacity:1; transform: none; }
}
.sr-hero::before {
    content: '';
    position: absolute; top: -1px; left: 8%; right: 8%; height: 2px;
    background: linear-gradient(90deg,
        transparent, rgba(99,102,241,0.9), rgba(6,182,212,0.7), rgba(139,92,246,0.6), transparent);
    border-radius: 100px;
    animation: srTopLine 4s ease-in-out infinite alternate;
}
@keyframes srTopLine {
    from { opacity: 0.6; left: 8%; right: 8%; }
    to   { opacity: 1.0; left: 4%; right: 4%; }
}
.sr-hero::after {
    content: '';
    position: absolute; inset: 0;
    background: url("data:image/svg+xml,%3Csvg width='60' height='60' xmlns='http://www.w3.org/2000/svg'%3E%3Cpath d='M0 60L60 0' stroke='rgba(255,255,255,0.015)' stroke-width='1'/%3E%3C/svg%3E") repeat;
    pointer-events: none;
}
.sr-hero-eyebrow {
    display: inline-flex; align-items: center; gap: 8px;
    background: rgba(99,102,241,0.10);
    border: 1px solid rgba(99,102,241,0.28);
    border-radius: 100px; padding: 6px 18px;
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.68rem; color: #a5b4fc; letter-spacing: 2.5px;
    text-transform: uppercase; margin-bottom: 20px;
    position: relative; z-index: 1;
    animation: srFadeUp 0.7s 0.1s both;
}
.sr-hero-dot {
    width: 7px; height: 7px; border-radius: 50%; background: #818cf8;
    box-shadow: 0 0 8px #6366f188;
    animation: srDot 1.6s ease-in-out infinite;
}
@keyframes srDot { 0%,100%{opacity:1;transform:scale(1);} 50%{opacity:0.3;transform:scale(0.7);} }
.sr-hero-title {
    font-family: 'Syne', sans-serif;
    font-size: clamp(2rem, 4vw, 3rem); font-weight: 800;
    line-height: 1.15; margin-bottom: 14px;
    background: linear-gradient(135deg,
        #ffffff 0%, #c7d2fe 25%, #818cf8 45%, #06b6d4 70%, #a5f3fc 100%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text;
    filter: drop-shadow(0 0 30px rgba(99,102,241,0.35));
    position: relative; z-index: 1;
    animation: srFadeUp 0.7s 0.15s both;
}
.sr-hero-sub {
    font-size: 1rem; font-weight: 400;
    color: rgba(255,255,255,0.5); line-height: 1.75;
    max-width: 600px; margin: 0 auto;
    position: relative; z-index: 1;
    animation: srFadeUp 0.7s 0.2s both;
}
.sr-hero-sub strong { color: rgba(165,180,252,0.9); font-weight: 600; }

/* ═══ FORMAT PILLS ═══ */
.sr-pills-row {
    display: flex; gap: 10px; flex-wrap: wrap;
    justify-content: center; margin: 20px 0 28px;
    animation: srFadeUp 0.7s 0.25s both;
}
.sr-pill {
    display: inline-flex; align-items: center; gap: 8px;
    padding: 9px 20px; border-radius: 50px;
    font-weight: 600; font-size: 0.82rem;
    border: 1px solid rgba(255,255,255,0.1);
    background: rgba(255,255,255,0.04);
    color: rgba(255,255,255,0.55);
    transition: all 0.3s cubic-bezier(0.16,1,0.3,1);
    cursor: default; letter-spacing: 0.3px;
}
.sr-pill:hover {
    background: rgba(99,102,241,0.14);
    border-color: rgba(99,102,241,0.45);
    color: #c7d2fe;
    transform: translateY(-2px);
    box-shadow: 0 8px 24px rgba(99,102,241,0.2);
}
.sr-pill-count {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.7rem; padding: 2px 8px;
    border-radius: 100px; background: rgba(99,102,241,0.2);
    color: #a5b4fc; letter-spacing: 1px;
}

/* ═══ UPLOAD ZONE ═══ */
.sr-dropzone {
    position: relative; overflow: hidden;
    background: rgba(10,14,30,0.7);
    border: 2px dashed rgba(99,102,241,0.35);
    border-radius: 22px; padding: 48px 32px;
    text-align: center; margin-bottom: 24px;
    transition: all 0.35s cubic-bezier(0.16,1,0.3,1);
    cursor: pointer;
}
.sr-dropzone::before {
    content: '';
    position: absolute; inset: 0;
    background: radial-gradient(ellipse at 50% 0%, rgba(99,102,241,0.08), transparent 70%);
    pointer-events: none;
}
.sr-dropzone:hover {
    border-color: rgba(99,102,241,0.7);
    background: rgba(99,102,241,0.06);
    transform: translateY(-2px);
    box-shadow: 0 20px 60px rgba(99,102,241,0.12);
}
.sr-dz-icon {
    font-size: 4rem; display: block; margin-bottom: 14px;
    animation: srIconFloat 3s ease-in-out infinite;
}
@keyframes srIconFloat {
    0%,100% { transform: translateY(0) rotate(-2deg); }
    50%      { transform: translateY(-10px) rotate(2deg); }
}
.sr-dz-title {
    font-family: 'Syne', sans-serif; font-size: 1.2rem; font-weight: 700;
    color: rgba(255,255,255,0.88); margin-bottom: 8px;
}
.sr-dz-sub {
    font-size: 0.82rem; color: rgba(255,255,255,0.32);
    line-height: 1.7; letter-spacing: 0.3px;
}
.sr-dz-formats {
    display: flex; gap: 6px; justify-content: center;
    flex-wrap: wrap; margin-top: 14px;
}
.sr-dz-fmtbadge {
    padding: 4px 12px; border-radius: 100px;
    background: rgba(99,102,241,0.1);
    border: 1px solid rgba(99,102,241,0.2);
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.65rem; color: #818cf8; letter-spacing: 1px;
}

/* ═══ STAT BAR ═══ */
.sr-statbar {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(90px, 1fr));
    gap: 10px; margin-bottom: 24px;
}
.sr-stat-tile {
    background: rgba(10,14,30,0.8);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 16px; padding: 16px 14px;
    text-align: center; position: relative; overflow: hidden;
    transition: all 0.3s ease;
}
.sr-stat-tile::before {
    content: '';
    position: absolute; top: 0; left: 0; right: 0; height: 2px;
    background: linear-gradient(90deg, var(--tc1,#6366f1), var(--tc2,#06b6d4));
    border-radius: 100px;
}
.sr-stat-tile:hover {
    border-color: rgba(99,102,241,0.3);
    transform: translateY(-3px);
    box-shadow: 0 12px 40px rgba(0,0,0,0.3);
}
.sr-stat-n {
    font-family: 'Syne', sans-serif; font-size: 1.6rem; font-weight: 800;
    background: linear-gradient(135deg, var(--tc1,#6366f1), var(--tc2,#06b6d4));
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; line-height: 1.1; margin-bottom: 4px;
}
.sr-stat-l {
    font-size: 0.62rem; color: rgba(255,255,255,0.3);
    letter-spacing: 2px; text-transform: uppercase;
    font-family: 'JetBrains Mono', monospace;
}

/* ═══ SECTION HEADER ═══ */
.sr-section-head {
    display: flex; align-items: center; gap: 10px;
    margin: 24px 0 14px;
}
.sr-section-line {
    flex: 1; height: 1px;
    background: linear-gradient(90deg, rgba(99,102,241,0.4), transparent);
}
.sr-section-label {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.65rem; letter-spacing: 3px;
    color: #818cf8; text-transform: uppercase;
    white-space: nowrap;
}

/* ═══ FUNCTION BUTTON GRID ═══ */
.sr-fnbtn {
    position: relative; overflow: hidden;
    background: rgba(10,14,30,0.7);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 16px; padding: 16px 12px;
    text-align: center; cursor: pointer;
    transition: all 0.3s cubic-bezier(0.16,1,0.3,1);
}
.sr-fnbtn::before {
    content: '';
    position: absolute; inset: 0;
    background: radial-gradient(ellipse at top, var(--fnc, rgba(99,102,241,0.12)), transparent 70%);
    opacity: 0; transition: opacity 0.3s;
}
.sr-fnbtn:hover { transform: translateY(-5px); border-color: var(--fnb, rgba(99,102,241,0.5)); }
.sr-fnbtn:hover::before { opacity: 1; }
.sr-fnbtn:hover { box-shadow: 0 12px 40px rgba(0,0,0,0.35), 0 0 0 1px var(--fnb, rgba(99,102,241,0.3)); }
.sr-fnbtn-icon { font-size: 2rem; display: block; margin-bottom: 8px; transition: transform 0.3s ease; }
.sr-fnbtn:hover .sr-fnbtn-icon { transform: scale(1.2) rotate(-5deg); }
.sr-fnbtn-label {
    font-size: 0.76rem; font-weight: 700;
    color: rgba(255,255,255,0.78); line-height: 1.4;
    font-family: 'Plus Jakarta Sans', sans-serif;
}
.sr-fnbtn-desc {
    font-size: 0.65rem; color: rgba(255,255,255,0.28);
    margin-top: 4px; line-height: 1.4;
}

/* ═══ Q&A CHAT ═══ */
.sr-chat-wrap { display: flex; flex-direction: column; gap: 4px; margin-bottom: 16px; }
.sr-chat-q {
    align-self: flex-end;
    background: rgba(255,255,255,0.06);
    border: 1px solid rgba(255,255,255,0.1);
    border-radius: 18px 18px 4px 18px;
    padding: 12px 18px; max-width: 85%;
    font-size: 0.9rem; color: rgba(255,255,255,0.85); line-height: 1.65;
    animation: srBubble 0.35s cubic-bezier(0.16,1,0.3,1) both;
}
.sr-chat-a {
    align-self: flex-start;
    background: linear-gradient(135deg, rgba(99,102,241,0.1), rgba(6,182,212,0.06));
    border: 1px solid rgba(99,102,241,0.22);
    border-left: 3px solid #6366f1;
    border-radius: 4px 18px 18px 18px;
    padding: 14px 18px; max-width: 92%;
    font-size: 0.9rem; color: rgba(255,255,255,0.82); line-height: 1.75;
    animation: srBubble 0.35s cubic-bezier(0.16,1,0.3,1) both;
}
@keyframes srBubble {
    from { opacity:0; transform: translateY(8px) scale(0.98); }
    to   { opacity:1; transform: none; }
}
.sr-chat-meta {
    font-family: 'JetBrains Mono', monospace; font-size: 0.6rem;
    color: rgba(255,255,255,0.22); padding: 0 6px;
    letter-spacing: 1px;
}

/* ═══ RESULT DISPLAY ═══ */
.sr-result {
    position: relative; overflow: hidden;
    background: rgba(10,14,30,0.85);
    border: 1px solid rgba(99,102,241,0.22);
    border-radius: 22px; padding: 28px 26px;
    margin-top: 22px; backdrop-filter: blur(20px);
    animation: srFadeUp 0.5s cubic-bezier(0.16,1,0.3,1) both;
}
.sr-result::before {
    content: '';
    position: absolute; top: -1px; left: 5%; right: 5%; height: 1px;
    background: linear-gradient(90deg, transparent, #6366f1, #06b6d4, transparent);
    border-radius: 100px;
}
.sr-result-hd {
    display: flex; align-items: center; gap: 12px;
    margin-bottom: 18px; padding-bottom: 14px;
    border-bottom: 1px solid rgba(255,255,255,0.05);
}
.sr-result-ico { font-size: 1.6rem; }
.sr-result-ttl {
    font-family: 'Syne', sans-serif; font-size: 1.1rem; font-weight: 700;
    color: #fff; letter-spacing: 0.3px; flex: 1;
}
.sr-result-tag {
    font-family: 'JetBrains Mono', monospace; font-size: 0.6rem;
    color: #6366f1; background: rgba(99,102,241,0.1);
    border: 1px solid rgba(99,102,241,0.2);
    border-radius: 100px; padding: 3px 10px; letter-spacing: 2px;
}
.sr-result-body {
    font-size: 0.93rem; color: rgba(255,255,255,0.72);
    line-height: 1.9; white-space: pre-wrap;
    font-family: 'Plus Jakarta Sans', sans-serif;
}

/* ═══ PAGE VIEWER ═══ */
.sr-page-card {
    background: rgba(10,14,30,0.6);
    border: 1px solid rgba(255,255,255,0.06);
    border-radius: 16px; padding: 18px 20px;
    margin-bottom: 14px;
    transition: all 0.25s ease;
    position: relative;
}
.sr-page-card:hover {
    border-color: rgba(99,102,241,0.35);
    transform: translateX(4px);
    box-shadow: -4px 0 20px rgba(99,102,241,0.1);
}
.sr-page-num {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.62rem; color: #818cf8;
    letter-spacing: 3px; margin-bottom: 10px;
    text-transform: uppercase; display: flex; align-items: center; gap: 8px;
}
.sr-page-num::after {
    content: ''; flex: 1; height: 1px;
    background: linear-gradient(90deg, rgba(99,102,241,0.3), transparent);
}
.sr-page-text {
    font-size: 0.85rem; color: rgba(255,255,255,0.58);
    line-height: 1.8; font-family: 'Plus Jakarta Sans', sans-serif;
}

/* ═══ SLIDE CARD ═══ */
.sr-slide-card {
    position: relative; overflow: hidden;
    background: rgba(10,14,30,0.7);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 18px; padding: 22px 22px;
    margin-bottom: 16px;
    transition: all 0.3s cubic-bezier(0.16,1,0.3,1);
}
.sr-slide-card::before {
    content: '';
    position: absolute; left: 0; top: 0; bottom: 0; width: 3px;
    background: linear-gradient(180deg, #6366f1, #06b6d4);
    border-radius: 3px 0 0 3px;
}
.sr-slide-card:hover {
    border-color: rgba(99,102,241,0.3);
    transform: translateY(-3px);
    box-shadow: 0 16px 50px rgba(0,0,0,0.35);
}
.sr-slide-num {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.62rem; color: #818cf8;
    letter-spacing: 3px; text-transform: uppercase; margin-bottom: 8px;
}
.sr-slide-title {
    font-family: 'Syne', sans-serif; font-size: 1.05rem; font-weight: 700;
    color: rgba(255,255,255,0.92); margin-bottom: 10px;
}
.sr-slide-body {
    font-size: 0.84rem; color: rgba(255,255,255,0.52); line-height: 1.75;
}
.sr-slide-notes {
    background: rgba(6,182,212,0.06);
    border: 1px solid rgba(6,182,212,0.18);
    border-radius: 10px; padding: 10px 14px;
    font-size: 0.78rem; color: rgba(6,182,212,0.75);
    margin-top: 12px; line-height: 1.65;
}
.sr-slide-num-badge {
    position: absolute; top: 18px; right: 18px;
    background: rgba(99,102,241,0.12);
    border: 1px solid rgba(99,102,241,0.2);
    border-radius: 8px; padding: 4px 10px;
    font-family: 'Syne', sans-serif; font-size: 0.85rem; font-weight: 800;
    color: #818cf8;
}

/* ═══ FLASHCARD ═══ */
.sr-flashcard {
    position: relative; overflow: hidden;
    background: linear-gradient(135deg,
        rgba(99,102,241,0.13) 0%,
        rgba(6,182,212,0.08) 100%);
    border: 1px solid rgba(99,102,241,0.28);
    border-radius: 20px; padding: 30px 26px;
    text-align: center; margin-bottom: 14px;
    min-height: 150px;
    display: flex; flex-direction: column;
    align-items: center; justify-content: center;
    transition: all 0.3s ease;
    animation: srFadeUp 0.4s cubic-bezier(0.16,1,0.3,1) both;
}
.sr-flashcard::before {
    content: '';
    position: absolute; top: -1px; left: 15%; right: 15%; height: 1px;
    background: linear-gradient(90deg, transparent, rgba(99,102,241,0.8), transparent);
}
.sr-flashcard:hover {
    transform: translateY(-4px) scale(1.01);
    box-shadow: 0 20px 60px rgba(99,102,241,0.2);
    border-color: rgba(99,102,241,0.5);
}
.sr-fc-label {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.62rem; letter-spacing: 3px; color: #818cf8;
    text-transform: uppercase; margin-bottom: 14px;
    background: rgba(99,102,241,0.1); border: 1px solid rgba(99,102,241,0.2);
    border-radius: 100px; padding: 3px 12px; display: inline-block;
}
.sr-fc-text {
    font-size: 1.05rem; color: rgba(255,255,255,0.9);
    line-height: 1.7; font-weight: 500;
    font-family: 'Plus Jakarta Sans', sans-serif;
}

/* ═══ QUIZ CARD ═══ */
.sr-quiz-card {
    background: rgba(10,14,30,0.75);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 18px; padding: 22px 22px;
    margin-bottom: 18px;
    animation: srFadeUp 0.4s cubic-bezier(0.16,1,0.3,1) both;
}
.sr-quiz-q {
    font-size: 0.97rem; font-weight: 700;
    color: rgba(255,255,255,0.92); margin-bottom: 16px;
    line-height: 1.65; font-family: 'Plus Jakarta Sans', sans-serif;
}
.sr-quiz-qnum {
    display: inline-block; margin-right: 8px;
    background: rgba(99,102,241,0.15);
    border: 1px solid rgba(99,102,241,0.25);
    border-radius: 8px; padding: 1px 8px;
    font-family: 'Syne', sans-serif; font-size: 0.8rem; font-weight: 800;
    color: #818cf8;
}
.sr-quiz-opt {
    background: rgba(255,255,255,0.03);
    border: 1px solid rgba(255,255,255,0.08);
    border-radius: 12px; padding: 11px 16px;
    font-size: 0.87rem; color: rgba(255,255,255,0.6);
    margin-bottom: 8px; cursor: pointer;
    transition: all 0.22s ease;
    font-family: 'Plus Jakarta Sans', sans-serif;
}
.sr-quiz-opt:hover {
    border-color: rgba(99,102,241,0.4);
    background: rgba(99,102,241,0.06);
    color: rgba(255,255,255,0.92);
    transform: translateX(4px);
}
.sr-quiz-ans {
    background: rgba(16,185,129,0.1);
    border: 1px solid rgba(16,185,129,0.3);
    border-radius: 12px; padding: 12px 16px;
    font-size: 0.85rem; color: #6ee7b7;
    margin-top: 12px; line-height: 1.6;
    font-family: 'Plus Jakarta Sans', sans-serif;
}

/* ═══ KEYWORD CHIPS ═══ */
.sr-kw-wrap { line-height: 2.4; }
.sr-kw {
    display: inline-block;
    border-radius: 100px; padding: 5px 16px;
    font-size: 0.78rem; font-weight: 600;
    margin: 4px; transition: all 0.2s ease; cursor: default;
    font-family: 'Plus Jakarta Sans', sans-serif;
}
.sr-kw:hover { transform: translateY(-2px) scale(1.05); }
.sr-kw-a { background: rgba(99,102,241,0.15); border: 1px solid rgba(99,102,241,0.3); color: #c7d2fe; }
.sr-kw-b { background: rgba(6,182,212,0.12);  border: 1px solid rgba(6,182,212,0.25);  color: #a5f3fc; }
.sr-kw-c { background: rgba(139,92,246,0.13); border: 1px solid rgba(139,92,246,0.28); color: #ddd6fe; }
.sr-kw-d { background: rgba(245,158,11,0.12); border: 1px solid rgba(245,158,11,0.25); color: #fde68a; }
.sr-kw-e { background: rgba(16,185,129,0.11); border: 1px solid rgba(16,185,129,0.25); color: #a7f3d0; }

/* ═══ TOC ═══ */
.sr-toc-row {
    display: flex; align-items: center; gap: 10px;
    padding: 9px 12px; border-radius: 10px;
    transition: all 0.2s;
}
.sr-toc-row:hover { background: rgba(99,102,241,0.06); }
.sr-toc-pip {
    width: 6px; height: 6px; border-radius: 50%;
    background: linear-gradient(135deg, #6366f1, #06b6d4); flex-shrink: 0;
}
.sr-toc-txt {
    font-size: 0.86rem; color: rgba(255,255,255,0.62); line-height: 1.5;
    font-family: 'Plus Jakarta Sans', sans-serif;
}

/* ═══ VIDEO FRAME GRID ═══ */
.sr-frame-wrap {
    border-radius: 16px; overflow: hidden;
    border: 1px solid rgba(255,255,255,0.08);
    background: rgba(10,14,30,0.6);
    transition: all 0.3s ease;
}
.sr-frame-wrap:hover {
    border-color: rgba(99,102,241,0.35);
    transform: translateY(-4px);
    box-shadow: 0 16px 50px rgba(0,0,0,0.4);
}
.sr-frame-ts {
    font-family: 'JetBrains Mono', monospace;
    font-size: 0.65rem; letter-spacing: 2px; color: #818cf8;
    padding: 10px 14px 6px; text-transform: uppercase;
}
.sr-frame-lbl {
    font-size: 0.8rem; color: rgba(255,255,255,0.55);
    padding: 0 14px 12px; line-height: 1.6;
}

/* ═══ DIVIDER ═══ */
.sr-div {
    height: 1px; margin: 22px 0;
    background: linear-gradient(90deg, transparent, rgba(99,102,241,0.35), rgba(6,182,212,0.2), transparent);
}

/* ═══ IMAGE PREVIEW ═══ */
.sr-img-frame {
    border-radius: 18px; overflow: hidden;
    border: 1px solid rgba(255,255,255,0.08);
    box-shadow: 0 24px 80px rgba(0,0,0,0.5);
    transition: all 0.35s ease;
}
.sr-img-frame:hover { transform: scale(1.01); box-shadow: 0 32px 100px rgba(99,102,241,0.2); }

/* ═══ SCROLLBAR ═══ */
.sr-root ::-webkit-scrollbar { width: 5px; }
.sr-root ::-webkit-scrollbar-track { background: transparent; }
.sr-root ::-webkit-scrollbar-thumb { background: rgba(99,102,241,0.3); border-radius: 10px; }
.sr-root ::-webkit-scrollbar-thumb:hover { background: rgba(99,102,241,0.6); }
</style>
"""


# ═══════════════════════════════════════════════════════════
# AI FUNCTION REGISTRY
# ═══════════════════════════════════════════════════════════

PDF_FUNCTIONS = [
    {"id": "summarize",    "icon": "📋", "label": "Summarize",        "desc": "Executive summary of the document",     "color": "#6366f1"},
    {"id": "qa",           "icon": "💬", "label": "Q&A Chat",         "desc": "Ask anything about the document",       "color": "#8b5cf6"},
    {"id": "flashcards",   "icon": "🃏", "label": "Flashcards",       "desc": "Auto-generate study flashcards",        "color": "#06b6d4"},
    {"id": "quiz",         "icon": "🧠", "label": "Smart Quiz",       "desc": "Multiple-choice quiz from content",     "color": "#10b981"},
    {"id": "keywords",     "icon": "🔑", "label": "Keywords",         "desc": "Top terms and key concepts",            "color": "#f59e0b"},
    {"id": "translate",    "icon": "🌍", "label": "Translate",        "desc": "Translate entire document",             "color": "#ec4899"},
    {"id": "explain5",     "icon": "👶", "label": "Explain Like 5",   "desc": "Simple plain-English explanation",      "color": "#f97316"},
    {"id": "chapters",     "icon": "📖", "label": "Chapter Analysis", "desc": "Section-by-section breakdown",         "color": "#6366f1"},
    {"id": "actions",      "icon": "✅", "label": "Action Items",     "desc": "Extract tasks and to-dos",             "color": "#10b981"},
    {"id": "critique",     "icon": "🎯", "label": "Critical Review",  "desc": "Strengths, weaknesses, improvements",  "color": "#ef4444"},
    {"id": "mindmap",      "icon": "🗺️", "label": "Mind Map Text",   "desc": "Hierarchical concept map outline",     "color": "#8b5cf6"},
    {"id": "timeline",     "icon": "⏳", "label": "Timeline",         "desc": "Chronological events extracted",        "color": "#f59e0b"},
    {"id": "simplify",     "icon": "✨", "label": "Simplify",         "desc": "Plain, simple language rewrite",        "color": "#06b6d4"},
    {"id": "toc",          "icon": "📑", "label": "Table of Contents","desc": "Auto-generated structure/TOC",          "color": "#a78bfa"},
    {"id": "sentiment",    "icon": "😊", "label": "Tone & Sentiment", "desc": "Emotional tone analysis",              "color": "#ec4899"},
    {"id": "citations",    "icon": "📚", "label": "Extract Citations","desc": "All references and citations listed",   "color": "#f97316"},
]

IMG_FUNCTIONS = [
    {"id": "describe",     "icon": "🏷️", "label": "Describe Image",  "desc": "Full AI image description",           "color": "#6366f1"},
    {"id": "ocr",          "icon": "🔤", "label": "Extract Text",     "desc": "OCR — read all text in image",         "color": "#8b5cf6"},
    {"id": "diagram",      "icon": "📊", "label": "Explain Diagram",  "desc": "Charts, diagrams, infographics",       "color": "#06b6d4"},
    {"id": "handwriting",  "icon": "✍️", "label": "Read Handwriting", "desc": "Decode handwritten notes",            "color": "#10b981"},
    {"id": "formula",      "icon": "🧮", "label": "Solve Formula",    "desc": "Math / science formula extraction",   "color": "#f59e0b"},
    {"id": "objects",      "icon": "📦", "label": "Object List",      "desc": "Enumerate everything in the image",    "color": "#ec4899"},
    {"id": "sentiment_img","icon": "🎭", "label": "Mood & Sentiment", "desc": "Emotional tone / scene mood",          "color": "#f97316"},
    {"id": "alttext",      "icon": "♿", "label": "Alt Text",         "desc": "Accessibility description",            "color": "#6366f1"},
    {"id": "artistic",     "icon": "🎨", "label": "Artistic Analysis","desc": "Style, composition, color palette",    "color": "#8b5cf6"},
    {"id": "chart_data",   "icon": "📈", "label": "Chart to Data",    "desc": "Extract numbers from charts/graphs",  "color": "#06b6d4"},
    {"id": "qa_img",       "icon": "💬", "label": "Ask About Image",  "desc": "Free-form Q&A on the visual",         "color": "#10b981"},
    {"id": "translate_img","icon": "🌍", "label": "Translate Text",   "desc": "Translate text found in image",       "color": "#ec4899"},
]

PPT_FUNCTIONS = [
    {"id": "exec_summary", "icon": "📋", "label": "Exec Summary",     "desc": "High-level presentation overview",    "color": "#6366f1"},
    {"id": "slide_notes",  "icon": "🎤", "label": "Speaker Notes",    "desc": "Generate notes for each slide",       "color": "#8b5cf6"},
    {"id": "quiz_ppt",     "icon": "🧠", "label": "Quiz From Slides", "desc": "Q&A quiz based on slide content",     "color": "#06b6d4"},
    {"id": "slide_summary","icon": "📖", "label": "Slide Summaries",  "desc": "Individual summary per slide",        "color": "#10b981"},
    {"id": "audience_qa",  "icon": "🙋", "label": "Audience Q&A",    "desc": "Likely questions audience will ask",   "color": "#f59e0b"},
    {"id": "critique_ppt", "icon": "🎯", "label": "Presentation Tips","desc": "Structure & delivery improvements",   "color": "#ef4444"},
    {"id": "story_arc",    "icon": "📢", "label": "Story Arc",        "desc": "Narrative structure analysis",        "color": "#ec4899"},
    {"id": "keywords_ppt", "icon": "🔑", "label": "Key Themes",       "desc": "Main topics and recurring ideas",     "color": "#f97316"},
    {"id": "translate_ppt","icon": "🌍", "label": "Translate PPT",    "desc": "Translate slide content",             "color": "#a78bfa"},
    {"id": "flashcards_ppt","icon":"🃏", "label": "Flashcards",       "desc": "Study cards from slide content",      "color": "#06b6d4"},
]

VID_FUNCTIONS = [
    {"id": "frame_desc",   "icon": "🎬", "label": "Frame Analysis",   "desc": "Describe sampled video frames",       "color": "#6366f1"},
    {"id": "key_moments",  "icon": "⭐", "label": "Key Moments",      "desc": "Identify important scenes",           "color": "#f59e0b"},
    {"id": "action_items_v","icon":"✅", "label": "Action Items",     "desc": "Tasks / instructions from video",     "color": "#10b981"},
    {"id": "topic_timeline","icon":"⏳", "label": "Topic Timeline",   "desc": "When each topic appears",              "color": "#8b5cf6"},
    {"id": "script",       "icon": "📝", "label": "Generate Script",  "desc": "Estimated script from visual context","color": "#06b6d4"},
    {"id": "qa_vid",       "icon": "💬", "label": "Ask About Video",  "desc": "Q&A based on extracted frames",       "color": "#ec4899"},
]


# ═══════════════════════════════════════════════════════════
# AI FUNCTION IMPLEMENTATIONS
# ═══════════════════════════════════════════════════════════

def _run_pdf_fn(fn_id: str, text: str, extra: str = "") -> str:
    sys = "You are an elite AI document analyst. Be thorough, structured, and insightful."
    prompts = {
        "summarize": f"Write a comprehensive executive summary of this document. Include: main topic, key arguments, important findings, and conclusions. Use headers and bullet points.\n\nDocument:\n{text[:25000]}",
        "flashcards": f"Create 15 study flashcards from this document. Format as:\nQ: [question]\nA: [answer]\n\nSeparate each with ---\n\nDocument:\n{text[:20000]}",
        "quiz": f"Create 10 multiple-choice questions from this document. Format each as:\nQ: [question]\nA) [option]\nB) [option]\nC) [option]\nD) [option]\nANSWER: [letter] — [explanation]\n\nSeparate each with ---\n\nDocument:\n{text[:20000]}",
        "keywords": f"Extract the 30 most important keywords, terms, and concepts from this document. Group them by theme. Bold the most critical ones.\n\nDocument:\n{text[:20000]}",
        "translate": f"Translate this document into {extra or 'Hindi'}. Maintain the original structure and formatting.\n\nDocument:\n{text[:15000]}",
        "explain5": f"Explain this entire document as if talking to a 10-year-old. Use very simple words, fun examples, and short sentences. Make it engaging!\n\nDocument:\n{text[:15000]}",
        "chapters": f"Analyze this document section by section. For each major section: title, 3-sentence summary, key points. Format clearly.\n\nDocument:\n{text[:25000]}",
        "actions": f"Extract ALL action items, tasks, recommendations, and to-dos from this document. Format as a prioritized checklist. Include page references if visible.\n\nDocument:\n{text[:20000]}",
        "critique": f"Provide a critical academic review of this document. Cover: strengths, weaknesses, logical gaps, missing context, quality of evidence, and improvement suggestions.\n\nDocument:\n{text[:20000]}",
        "mindmap": f"Create a detailed hierarchical mind map outline of this document. Use indentation levels (Main → Sub → Detail). Include all key concepts.\n\nDocument:\n{text[:20000]}",
        "timeline": f"Extract all dates, events, and time references from this document and arrange them as a chronological timeline. Include year, event, and significance.\n\nDocument:\n{text[:20000]}",
        "simplify": f"Rewrite this document in simple, clear English. Remove jargon, use short sentences, keep all key information, add clear headings.\n\nDocument:\n{text[:15000]}",
        "toc": f"Generate a detailed Table of Contents for this document. Identify all major sections, subsections, and key topics. Number them hierarchically.\n\nDocument:\n{text[:20000]}",
        "sentiment": f"Analyze the tone, sentiment, and writing style of this document. Cover: overall tone (formal/informal/academic/persuasive), sentiment (positive/neutral/negative), emotional register, and rhetorical devices used.\n\nDocument:\n{text[:15000]}",
        "citations": f"Extract all citations, references, bibliographic entries, footnotes, and sources mentioned in this document. Format them as a numbered reference list.\n\nDocument:\n{text[:20000]}",
    }
    p = prompts.get(fn_id, f"Analyze this document:\n{text[:15000]}")
    if fn_id == "qa":
        p = f"Answer this question about the document:\n\nQuestion: {extra}\n\nDocument:\n{text[:20000]}"
    return _call_ai(p, sys, max_tokens=3000, temperature=0.4)


def _run_img_fn(fn_id: str, b64: str, mime: str, extra: str = "") -> str:
    prompts = {
        "describe": "Describe this image in rich detail. Include: subjects, objects, setting, colors, lighting, composition, mood, and any text visible.",
        "ocr": "Extract ALL text visible in this image exactly as it appears. Preserve formatting, line breaks, and layout. If no text, say so.",
        "diagram": "This appears to be a diagram/chart/infographic. Explain it completely: what it shows, the data or concepts represented, relationships, and key takeaways.",
        "handwriting": "Read and transcribe ALL handwritten text in this image exactly as written. If any part is unclear, note it with [unclear].",
        "formula": "Extract all mathematical, chemical, or scientific formulas from this image. Write them out clearly and explain what each formula represents and how to use it.",
        "objects": "List every single object, element, person, animal, shape, and item visible in this image. Be exhaustive and specific.",
        "sentiment_img": "Analyze the mood, emotion, and sentiment conveyed by this image. Describe the atmosphere, emotional tone, and what feelings it evokes.",
        "alttext": "Write professional accessibility alt-text for this image. Be specific, descriptive, and concise (2-3 sentences). Follow WCAG guidelines.",
        "artistic": "Analyze this image artistically. Cover: visual style, composition techniques, color palette, lighting, artistic movement/influence, technical skill, symbolism, and overall aesthetic quality.",
        "chart_data": "This is a chart or graph. Extract the exact data: axes labels, data series names, approximate values for all data points, trends, and key statistics visible.",
        "translate_img": "Extract all text from this image and translate it into English. Show: original text, then translation, maintaining structure.",
    }
    p = prompts.get(fn_id, "Analyze this image thoroughly.")
    if fn_id == "qa_img":
        p = f"Look at this image carefully and answer: {extra or 'What is shown in this image?'}"
    return _call_vision(p, b64, mime, max_tokens=2500)


def _run_ppt_fn(fn_id: str, slides: List[dict], full_text: str, extra: str = "") -> str:
    sys = "You are an expert presentation coach and content strategist."
    ctx = full_text[:20000]
    slides_str = "\n\n".join(
        f"Slide {s['num']}: {s['title']}\n{s['body']}" for s in slides[:30]
    )
    prompts = {
        "exec_summary": f"Write a concise executive summary of this presentation. Include: main message, key points per section, and conclusion.\n\n{slides_str}",
        "slide_notes": f"Generate professional speaker notes for each slide — 3-4 sentences each, including talking points, transitions, and emphasis cues.\n\n{slides_str}",
        "quiz_ppt": f"Create 10 multiple-choice quiz questions based on this presentation content.\n\n{slides_str}",
        "slide_summary": f"Write a 2-sentence summary for each slide, capturing the core message.\n\n{slides_str}",
        "audience_qa": f"Generate 15 insightful questions that an engaged audience would likely ask after this presentation.\n\n{slides_str}",
        "critique_ppt": f"Critically review this presentation. Cover: structure, clarity, flow, slide design best practices, content gaps, and specific improvement suggestions per slide.\n\n{slides_str}",
        "story_arc": f"Analyze the narrative/story arc of this presentation. Describe the opening hook, tension-building, key arguments, and resolution. Rate the storytelling 1-10 with reasoning.\n\n{slides_str}",
        "keywords_ppt": f"Extract the main themes, key terms, and recurring concepts across all slides. Group them by importance.\n\n{slides_str}",
        "translate_ppt": f"Translate all slide content into {extra or 'Hindi'}. Maintain slide structure.\n\n{slides_str}",
        "flashcards_ppt": f"Create 15 study flashcards from the presentation content.\nQ: [question]\nA: [answer]\nSeparate with ---\n\n{slides_str}",
    }
    p = prompts.get(fn_id, f"Analyze this presentation:\n{slides_str}")
    return _call_ai(p, sys, max_tokens=3000, temperature=0.4)


def _run_vid_fn(fn_id: str, frames: List[Tuple[str, str, float]],
                extra: str = "") -> str:
    """Analyze video by processing key frames with vision AI."""
    if not frames:
        return "No frames could be extracted from the video."

    # Describe each frame
    frame_descs = []
    for i, (b64, mime, ts) in enumerate(frames):
        mins = int(ts // 60)
        secs = int(ts % 60)
        desc = _call_vision(
            f"Describe this video frame (timestamp {mins}:{secs:02d}) in detail.",
            b64, mime, max_tokens=400
        )
        frame_descs.append(f"[{mins}:{secs:02d}] {desc}")

    combined = "\n\n".join(frame_descs)
    sys = "You are analyzing a video based on extracted frame descriptions."

    prompts = {
        "frame_desc": combined,
        "key_moments": f"Based on these frame descriptions, identify the key moments, pivotal scenes, and most important events in the video. List with timestamps.\n\n{combined}",
        "action_items_v": f"Extract all action items, instructions, and tasks shown or implied in this video.\n\n{combined}",
        "topic_timeline": f"Create a timeline of topics covered throughout this video based on the frames. Format: [timestamp] - Topic/Content\n\n{combined}",
        "script": f"Based on the video frames, write an estimated narration/script for what likely appears in this video. Be creative but grounded in what's visible.\n\n{combined}",
        "qa_vid": f"Answer this question about the video: {extra or 'What is this video about?'}\n\nFrame descriptions:\n{combined}",
    }
    if fn_id == "frame_desc":
        return combined  # already done above
    return _call_ai(prompts.get(fn_id, combined), sys, max_tokens=2500)


# ═══════════════════════════════════════════════════════════
# RENDER HELPERS
# ═══════════════════════════════════════════════════════════

def _render_fn_grid(functions: List[dict], key_prefix: str) -> Optional[str]:
    """Render clickable function grid. Returns selected fn_id or None."""
    st.markdown('<div class="sr-fn-grid">', unsafe_allow_html=True)
    cols = st.columns(4)
    chosen = None
    for i, fn in enumerate(functions):
        with cols[i % 4]:
            if st.button(
                f"{fn['icon']} {fn['label']}",
                key=f"{key_prefix}_fn_{fn['id']}",
                use_container_width=True,
                help=fn["desc"],
            ):
                chosen = fn["id"]
    st.markdown('</div>', unsafe_allow_html=True)
    return chosen


def _render_result(icon: str, title: str, body: str,
                   allow_download: bool = True, filename: str = "result.txt"):
    safe_body = body.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    st.markdown(f"""
    <div class="sr-result">
        <div class="sr-result-hd">
            <span class="sr-result-ico">{icon}</span>
            <span class="sr-result-ttl">{title}</span>
            <span class="sr-result-tag">AI RESULT</span>
        </div>
        <div class="sr-result-body">{safe_body}</div>
    </div>
    """, unsafe_allow_html=True)
    if allow_download and body:
        st.download_button(
            "📥 Download",
            data=body, file_name=filename, mime="text/plain",
            key=f"dl_{filename}_{random.randint(1000,9999)}",
        )


def _render_pdf_stats(meta: dict):
    pages  = meta.get("page_count", "?")
    words  = meta.get("word_count", 0)
    mins   = meta.get("reading_min", "?")
    title  = (meta.get("title", "") or "Unknown")[:20]
    author = (meta.get("author", "") or "Unknown")[:16]
    colors = [
        ("--tc1:#6366f1", "--tc2:#8b5cf6"),
        ("--tc1:#06b6d4", "--tc2:#0ea5e9"),
        ("--tc1:#10b981", "--tc2:#34d399"),
        ("--tc1:#f59e0b", "--tc2:#fb923c"),
        ("--tc1:#ec4899", "--tc2:#a78bfa"),
    ]
    tiles = [
        (pages, "PAGES"),
        (f"{words:,}" if isinstance(words, int) else words, "WORDS"),
        (f"{mins}m", "READ TIME"),
        (title, "TITLE"),
        (author, "AUTHOR"),
    ]
    html = '<div class="sr-statbar">'
    for i, (val, lbl) in enumerate(tiles):
        c1, c2 = colors[i % len(colors)]
        html += f'<div class="sr-stat-tile" style="{c1};{c2}">'
        html += f'<div class="sr-stat-n">{val}</div>'
        html += f'<div class="sr-stat-l">{lbl}</div>'
        html += '</div>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)


def _render_flashcards(raw: str):
    cards = [c.strip() for c in raw.split("---") if c.strip()]
    if not cards:
        _render_result("🃏", "Flashcards", raw)
        return
    st.markdown(f"""
    <div class="sr-section-head">
        <span class="sr-section-label">🃏 {len(cards)} flashcards generated</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)
    for i, card in enumerate(cards):
        lines = card.strip().split("\n")
        q_lines = [l for l in lines if l.startswith("Q:")]
        a_lines = [l for l in lines if l.startswith("A:")]
        q_text = q_lines[0][2:].strip() if q_lines else card
        a_text = a_lines[0][2:].strip() if a_lines else ""
        st.markdown(f"""
        <div class="sr-flashcard" style="animation-delay:{i*0.05}s">
            <div class="sr-fc-label">Card {i+1} of {len(cards)}</div>
            <div class="sr-fc-text">{q_text}</div>
        </div>
        """, unsafe_allow_html=True)
        if a_text:
            with st.expander(f"💡 Reveal Answer #{i+1}"):
                st.markdown(f"""
                <div style="padding:14px 18px;background:rgba(16,185,129,0.08);
                    border:1px solid rgba(16,185,129,0.25);border-radius:14px;
                    font-size:0.92rem;color:#6ee7b7;line-height:1.7;">
                    {a_text}
                </div>
                """, unsafe_allow_html=True)


def _render_quiz(raw: str):
    questions = [q.strip() for q in raw.split("---") if q.strip()]
    st.markdown(f"""
    <div class="sr-section-head">
        <span class="sr-section-label">🧠 {len(questions)} quiz questions</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)
    for i, qblock in enumerate(questions):
        lines = [l.strip() for l in qblock.strip().split("\n") if l.strip()]
        q_text = ""; opts = []; answer_line = ""
        for l in lines:
            if l.startswith("Q:"):         q_text = l[2:].strip()
            elif l[:2] in ("A)","B)","C)","D)"): opts.append(l)
            elif l.startswith("ANSWER:"): answer_line = l[7:].strip()
        st.markdown(f"""
        <div class="sr-quiz-card" style="animation-delay:{i*0.06}s">
            <div class="sr-quiz-q">
                <span class="sr-quiz-qnum">Q{i+1}</span>{q_text}
            </div>
            {''.join(f'<div class="sr-quiz-opt">{o}</div>' for o in opts)}
        </div>
        """, unsafe_allow_html=True)
        if answer_line:
            with st.expander(f"✅ Reveal Answer for Q{i+1}"):
                st.markdown(f'<div class="sr-quiz-ans">✅ {answer_line}</div>',
                            unsafe_allow_html=True)


_KW_CLASSES = ["sr-kw sr-kw-a", "sr-kw sr-kw-b", "sr-kw sr-kw-c", "sr-kw sr-kw-d", "sr-kw sr-kw-e"]

def _render_keywords(raw: str):
    seen = set(); chips_html = ""
    plain = [w.strip() for w in raw.replace("**", "").replace("`", "").split("\n") if w.strip()]
    idx = 0
    for chunk in plain[:50]:
        for part in re.split(r'[,;•\-–|/]', chunk):
            p = re.sub(r'^\d+\.\s*', '', part).strip()
            if 2 < len(p) < 55 and p not in seen and not p.startswith('#'):
                seen.add(p)
                cls = _KW_CLASSES[idx % len(_KW_CLASSES)]
                chips_html += f'<span class="{cls}">{p}</span>'
                idx += 1
    if chips_html:
        st.markdown(f'<div class="sr-kw-wrap">{chips_html}</div>', unsafe_allow_html=True)
    else:
        _render_result("🔑", "Keywords", raw)


# ═══════════════════════════════════════════════════════════
# FORMAT TABS
# ═══════════════════════════════════════════════════════════

def _render_pdf_tab():
    st.markdown('<div class="sr-upload-zone">'
                '<span class="sr-upload-icon">📄</span>'
                '<div class="sr-upload-title">Upload a PDF</div>'
                '<div class="sr-upload-sub">Research papers · Textbooks · Reports · Essays · Legal docs</div>'
                '</div>', unsafe_allow_html=True)

    pdf_file = st.file_uploader(
        "Choose PDF", type=["pdf"], key="sr_pdf_upload",
        label_visibility="collapsed"
    )

    if not pdf_file:
        return

    # Extract on first load
    cache_key = f"sr_pdf_{pdf_file.name}_{pdf_file.size}"
    if st.session_state.get("sr_pdf_cache_key") != cache_key:
        with st.spinner("🔍 Extracting text from PDF…"):
            raw = pdf_file.read()
            full_text, meta, page_texts = _extract_pdf_text(raw)
            st.session_state["sr_pdf_text"]       = full_text
            st.session_state["sr_pdf_meta"]       = meta
            st.session_state["sr_pdf_pages"]      = page_texts
            st.session_state["sr_pdf_cache_key"]  = cache_key
            st.session_state["sr_pdf_result"]     = {}
            st.session_state["sr_pdf_qa_history"] = []

    full_text  = st.session_state.get("sr_pdf_text", "")
    meta       = st.session_state.get("sr_pdf_meta", {})
    page_texts = st.session_state.get("sr_pdf_pages", [])

    if not full_text or full_text.startswith("Error"):
        st.error(f"Could not read PDF: {full_text}")
        return

    # Stats
    _render_pdf_stats(meta)

    # TOC if available
    toc = meta.get("toc", [])
    if toc:
        with st.expander("📑 Table of Contents", expanded=False):
            for lvl, title in toc[:25]:
                indent = (lvl - 1) * 18
                st.markdown(
                    f'<div class="sr-toc-row" style="padding-left:{indent}px">'
                    f'<div class="sr-toc-pip"></div>'
                    f'<div class="sr-toc-txt">{title}</div></div>',
                    unsafe_allow_html=True,
                )

    # Page preview
    with st.expander("📄 Page Viewer", expanded=False):
        total_pages = len(page_texts)
        pg = st.slider("Page", 1, max(1, total_pages), 1, key="sr_pdf_page_slider")
        if page_texts:
            st.markdown(f"""
            <div class="sr-page-card">
                <div class="sr-page-num">Page {pg} of {total_pages}</div>
                <div class="sr-page-text">{page_texts[pg-1][:2000]}</div>
            </div>
            """, unsafe_allow_html=True)

    st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
    st.markdown("""
    <div class="sr-section-head">
        <span class="sr-section-label">🧠 AI Functions</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)

    # Q&A chat panel
    st.markdown("""
    <div class="sr-section-head" style="margin-top:8px;">
        <span class="sr-section-label">💬 Document Q&amp;A Chat</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)
    qa_history = st.session_state.get("sr_pdf_qa_history", [])
    if qa_history:
        st.markdown('<div class="sr-chat-wrap">', unsafe_allow_html=True)
        for qa in qa_history:
            q_safe = qa["q"].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            a_safe = qa["a"].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            st.markdown(f'<div class="sr-chat-q">👤 {q_safe}</div>', unsafe_allow_html=True)
            st.markdown(f'<div class="sr-chat-a">🤖 {a_safe}</div>', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with st.form("sr_pdf_qa_form", clear_on_submit=True):
        q_input = st.text_input("Ask anything about this document…",
                                placeholder="What is the main argument? · Summarize page 3 · List all dates…",
                                label_visibility="collapsed", key="sr_pdf_qa_input")
        q_sub = st.form_submit_button("Ask →", use_container_width=True, type="primary")
    if q_sub and q_input.strip():
        with st.spinner("🤔 Thinking…"):
            ans = _run_pdf_fn("qa", full_text, extra=q_input)
        st.session_state["sr_pdf_qa_history"].append({"q": q_input, "a": ans})
        st.rerun()

    st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
    st.markdown("""
    <div class="sr-section-head">
        <span class="sr-section-label">⚡ AI Power Tools</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)

    # Function grid
    fn_map = {fn["id"]: fn for fn in PDF_FUNCTIONS if fn["id"] != "qa"}
    cols = st.columns(4)
    chosen_fn = None
    for i, fn in enumerate(list(fn_map.values())):
        with cols[i % 4]:
            if st.button(f"{fn['icon']} {fn['label']}", key=f"pdf_fn_{fn['id']}",
                         use_container_width=True, help=fn["desc"]):
                chosen_fn = fn["id"]

    # Translation language
    extra_input = ""
    if st.session_state.get("sr_pdf_expecting_translate"):
        extra_input = st.text_input("Target language:", value="Hindi", key="sr_pdf_lang")

    if chosen_fn:
        if chosen_fn == "translate":
            st.session_state["sr_pdf_expecting_translate"] = True
        else:
            st.session_state["sr_pdf_expecting_translate"] = False

        fn_info = fn_map.get(chosen_fn, {})
        with st.spinner(f"{fn_info.get('icon','⚙️')} Running {fn_info.get('label','AI')}…"):
            result = _run_pdf_fn(chosen_fn, full_text, extra=extra_input)
            st.session_state["sr_pdf_result"][chosen_fn] = (fn_info, result)

    # Display all results
    for fid, (fn_info, result) in st.session_state.get("sr_pdf_result", {}).items():
        st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
        if fid == "flashcards":
            _render_flashcards(result)
        elif fid == "quiz":
            _render_quiz(result)
        elif fid == "keywords":
            _render_keywords(result)
        else:
            _render_result(fn_info["icon"], fn_info["label"], result,
                           filename=f"pdf_{fid}.txt")


def _render_img_tab():
    st.markdown('<div class="sr-upload-zone">'
                '<span class="sr-upload-icon">🖼️</span>'
                '<div class="sr-upload-title">Upload an Image</div>'
                '<div class="sr-upload-sub">PNG · JPG · JPEG · WEBP · BMP · GIF</div>'
                '</div>', unsafe_allow_html=True)

    img_file = st.file_uploader(
        "Choose Image", type=["png", "jpg", "jpeg", "webp", "bmp", "gif"],
        key="sr_img_upload", label_visibility="collapsed"
    )

    if not img_file:
        return

    img_bytes = img_file.read()
    b64 = _bytes_to_b64(img_bytes)
    ext = img_file.name.split(".")[-1].lower()
    mime_map = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
                "webp": "image/webp", "bmp": "image/bmp", "gif": "image/gif"}
    mime = mime_map.get(ext, "image/jpeg")

    # Display image
    col_img, col_info = st.columns([1.4, 1])
    with col_img:
        st.image(img_bytes, use_container_width=True)
    with col_info:
        size_kb = len(img_bytes) / 1024
        st.markdown(f"""
        <div class="sr-stats-row" style="flex-direction:column;gap:8px;">
            <div class="sr-stat-chip">
                <div class="sr-stat-val">{size_kb:.0f}KB</div>
                <div class="sr-stat-lbl">File Size</div>
            </div>
            <div class="sr-stat-chip">
                <div class="sr-stat-val">{ext.upper()}</div>
                <div class="sr-stat-lbl">Format</div>
            </div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown('<div class="sr-divider"></div>', unsafe_allow_html=True)
    st.markdown("### 🤖 AI Functions")

    # Q&A special
    st.markdown("**💬 Ask About This Image**")
    with st.form("sr_img_qa_form", clear_on_submit=True):
        q_img = st.text_input("Your question about the image…",
                              placeholder="What is shown? · Read the text · Explain the chart…",
                              label_visibility="collapsed", key="sr_img_qa_in")
        q_img_sub = st.form_submit_button("Ask Vision AI →", use_container_width=True, type="primary")

    if q_img_sub and q_img.strip():
        with st.spinner("👁️ Analyzing image…"):
            ans = _run_img_fn("qa_img", b64, mime, extra=q_img)
        _render_result("💬", "Image Answer", ans, filename="img_answer.txt")

    st.markdown('<div class="sr-divider"></div>', unsafe_allow_html=True)

    # Function grid
    img_fns = [fn for fn in IMG_FUNCTIONS if fn["id"] != "qa_img"]
    cols = st.columns(4)
    chosen_fn = None
    for i, fn in enumerate(img_fns):
        with cols[i % 4]:
            if st.button(f"{fn['icon']} {fn['label']}", key=f"img_fn_{fn['id']}",
                         use_container_width=True, help=fn["desc"]):
                chosen_fn = fn["id"]

    if chosen_fn:
        fn_info = next((f for f in IMG_FUNCTIONS if f["id"] == chosen_fn), {})
        with st.spinner(f"{fn_info.get('icon','👁️')} Running {fn_info.get('label','Vision AI')}…"):
            result = _run_img_fn(chosen_fn, b64, mime)
        _render_result(fn_info.get("icon", "🏷️"),
                       fn_info.get("label", "Result"), result,
                       filename=f"img_{chosen_fn}.txt")


def _render_ppt_tab():
    st.markdown("""
    <div class="sr-dropzone">
        <span class="sr-dz-icon">📊</span>
        <div class="sr-dz-title">Drop your PowerPoint here</div>
        <div class="sr-dz-sub">Slide decks &middot; Lectures &middot; Business presentations &middot; Pitch decks</div>
        <div class="sr-dz-formats">
            <span class="sr-dz-fmtbadge">.PPT</span>
            <span class="sr-dz-fmtbadge">.PPTX</span>
        </div>
    </div>
    """, unsafe_allow_html=True)

    ppt_file = st.file_uploader(
        "Choose PPT/PPTX", type=["ppt", "pptx"], key="sr_ppt_upload",
        label_visibility="collapsed"
    )

    if not ppt_file:
        return

    cache_key = f"sr_ppt_{ppt_file.name}_{ppt_file.size}"
    if st.session_state.get("sr_ppt_cache_key") != cache_key:
        with st.spinner("📊 Reading PowerPoint…"):
            raw = ppt_file.read()
            slides, full_text = _extract_pptx(raw)
            st.session_state["sr_ppt_slides"]     = slides
            st.session_state["sr_ppt_text"]       = full_text
            st.session_state["sr_ppt_cache_key"]  = cache_key
            st.session_state["sr_ppt_result"]     = {}

    slides    = st.session_state.get("sr_ppt_slides", [])
    full_text = st.session_state.get("sr_ppt_text", "")

    if not slides:
        st.error(f"Could not read PPT: {full_text}")
        return

    # Stats with premium tiles
    n_slides = len(slides)
    n_words  = len(full_text.split())
    n_notes  = sum(1 for s in slides if s.get("notes"))
    colors   = [("--tc1:#6366f1","--tc2:#8b5cf6"),("--tc1:#06b6d4","--tc2:#0ea5e9"),("--tc1:#10b981","--tc2:#34d399")]
    tiles    = [(str(n_slides),"SLIDES"),(f"{n_words:,}","WORDS"),(str(n_notes),"WITH NOTES")]
    html = '<div class="sr-statbar">'
    for i,(val,lbl) in enumerate(tiles):
        c1,c2 = colors[i]
        html += f'<div class="sr-stat-tile" style="{c1};{c2}"><div class="sr-stat-n">{val}</div><div class="sr-stat-l">{lbl}</div></div>'
    html += '</div>'
    st.markdown(html, unsafe_allow_html=True)

    with st.expander("📊 Slide Viewer", expanded=True):
        for slide in slides[:20]:
            notes_html = ""
            if slide.get("notes"):
                notes_html = f'<div class="sr-slide-notes">📝 <b>Speaker Notes:</b> {slide["notes"][:300]}</div>'
            body_safe = slide["body"][:500].replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
            st.markdown(f"""
            <div class="sr-slide-card">
                <span class="sr-slide-num-badge">#{slide['num']}</span>
                <div class="sr-slide-num">SLIDE {slide['num']}</div>
                <div class="sr-slide-title">{slide['title']}</div>
                <div class="sr-slide-body">{body_safe}</div>
                {notes_html}
            </div>
            """, unsafe_allow_html=True)

    st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
    st.markdown("""
    <div class="sr-section-head">
        <span class="sr-section-label">⚡ AI Power Tools</span>
        <div class="sr-section-line"></div>
    </div>
    """, unsafe_allow_html=True)

    fn_map = {fn["id"]: fn for fn in PPT_FUNCTIONS}
    cols = st.columns(4)
    chosen_fn = None
    for i, fn in enumerate(list(fn_map.values())):
        with cols[i % 4]:
            if st.button(f"{fn['icon']} {fn['label']}", key=f"ppt_fn_{fn['id']}",
                         use_container_width=True, help=fn["desc"]):
                chosen_fn = fn["id"]

    extra_input = ""
    if chosen_fn == "translate_ppt":
        extra_input = st.text_input("Target language:", value="Hindi", key="sr_ppt_lang")

    if chosen_fn:
        fn_info = fn_map.get(chosen_fn, {})
        with st.spinner(f"{fn_info.get('icon','⚙️')} Running {fn_info.get('label','AI')}…"):
            result = _run_ppt_fn(chosen_fn, slides, full_text, extra=extra_input)
            st.session_state["sr_ppt_result"][chosen_fn] = (fn_info, result)

    for fid, (fn_info, result) in st.session_state.get("sr_ppt_result", {}).items():
        st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
        if fid == "quiz_ppt":
            _render_quiz(result)
        elif fid == "flashcards_ppt":
            _render_flashcards(result)
        else:
            _render_result(fn_info["icon"], fn_info["label"], result,
                           filename=f"ppt_{fid}.txt")


def _render_video_tab():
    st.markdown("""
    <div class="sr-dropzone">
        <span class="sr-dz-icon">🎬</span>
        <div class="sr-dz-title">Drop your Video here</div>
        <div class="sr-dz-sub">Frames are sampled &amp; analyzed by Vision AI &mdash; no audio needed</div>
        <div class="sr-dz-formats">
            <span class="sr-dz-fmtbadge">.MP4</span>
            <span class="sr-dz-fmtbadge">.AVI</span>
            <span class="sr-dz-fmtbadge">.MOV</span>
            <span class="sr-dz-fmtbadge">.MKV</span>
        </div>
    </div>
    """, unsafe_allow_html=True)

    if cv2 is None:
        st.warning("⚠️ OpenCV not installed. Install with: `pip install opencv-python`")

    vid_file = st.file_uploader(
        "Choose Video", type=["mp4", "avi", "mov", "mkv", "webm"],
        key="sr_vid_upload", label_visibility="collapsed"
    )

    if not vid_file:
        return

    cache_key = f"sr_vid_{vid_file.name}_{vid_file.size}"
    num_frames = st.slider("Frames to sample", 3, 10, 6, key="sr_vid_frames")

    if st.session_state.get("sr_vid_cache_key") != cache_key:
        st.session_state["sr_vid_frames_data"] = None
        st.session_state["sr_vid_cache_key"]   = cache_key
        st.session_state["sr_vid_result"]      = {}

    frames = st.session_state.get("sr_vid_frames_data")

    if frames is None:
        if st.button("🎬 Extract Frames (required first)", use_container_width=True, type="primary"):
            with st.spinner("🎞️ Extracting frames from video…"):
                raw = vid_file.read()
                frames = _extract_video_frames(raw, num_frames)
                st.session_state["sr_vid_frames_data"] = frames
                st.rerun()
    else:
        # Show frames as premium grid
        st.markdown(f"""
        <div class="sr-section-head">
            <span class="sr-section-label">🎬 {len(frames)} frames extracted</span>
            <div class="sr-section-line"></div>
        </div>
        """, unsafe_allow_html=True)
        f_cols = st.columns(3)
        for i, (b64_f, mime_f, ts) in enumerate(frames):
            mins = int(ts // 60); secs = int(ts % 60)
            with f_cols[i % 3]:
                img_data = base64.b64decode(b64_f)
                st.markdown('<div class="sr-frame-wrap">', unsafe_allow_html=True)
                st.image(img_data, use_container_width=True)
                st.markdown(f'<div class="sr-frame-ts">⏱ {mins}:{secs:02d}</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

        st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
        st.markdown("""
        <div class="sr-section-head">
            <span class="sr-section-label">🧠 AI Functions</span>
            <div class="sr-section-line"></div>
        </div>
        """, unsafe_allow_html=True)

        # Q&A special
        with st.form("sr_vid_qa_form", clear_on_submit=True):
            q_vid = st.text_input("Ask about this video…",
                                  placeholder="What happens in this video? · What is being taught?",
                                  label_visibility="collapsed", key="sr_vid_qa_in")
            q_vid_sub = st.form_submit_button("Ask AI →", use_container_width=True, type="primary")

        if q_vid_sub and q_vid.strip():
            with st.spinner("🎬 Analyzing video with AI…"):
                ans = _run_vid_fn("qa_vid", frames, extra=q_vid)
            _render_result("💬", "Video Answer", ans, filename="vid_answer.txt")

        fn_map = {fn["id"]: fn for fn in VID_FUNCTIONS if fn["id"] != "qa_vid"}
        cols = st.columns(3)
        chosen_fn = None
        for i, fn in enumerate(list(fn_map.values())):
            with cols[i % 3]:
                if st.button(f"{fn['icon']} {fn['label']}", key=f"vid_fn_{fn['id']}",
                             use_container_width=True, help=fn["desc"]):
                    chosen_fn = fn["id"]

        if chosen_fn:
            fn_info = fn_map.get(chosen_fn, {})
            with st.spinner(f"{fn_info.get('icon','🎬')} Running {fn_info.get('label','AI')} on frames…"):
                result = _run_vid_fn(chosen_fn, frames)
                st.session_state["sr_vid_result"][chosen_fn] = (fn_info, result)

        for fid, (fn_info, result) in st.session_state.get("sr_vid_result", {}).items():
            st.markdown('<div class="sr-div"></div>', unsafe_allow_html=True)
            _render_result(fn_info["icon"], fn_info["label"], result,
                           filename=f"vid_{fid}.txt")


def _render_media_tab():
    pass


# ═══════════════════════════════════════════════════════════
# MAIN RENDER
# ═══════════════════════════════════════════════════════════

def render_smart_reader():
    """Main entry point — render the AI Smart Reader page."""

    st.markdown(_CSS, unsafe_allow_html=True)

    # Init state
    for k, v in {
        "sr_pdf_result": {}, "sr_pdf_qa_history": [],
        "sr_ppt_result": {}, "sr_vid_result": {},
        "sr_pdf_expecting_translate": False,
    }.items():
        if k not in st.session_state:
            st.session_state[k] = v

    # Back button
    if st.button("← Back", key="sr_back"):
        st.session_state.app_mode = "chat"
        st.rerun()

    # Hero
    st.markdown("""
    <div class="sr-hero">
        <div style="margin-bottom:10px;">
            <span class="sr-badge">
                <span class="sr-badge-dot"></span>
                AI SMART READER &nbsp;·&nbsp; 4 FORMATS &nbsp;·&nbsp; 40+ AI FUNCTIONS
            </span>
        </div>
        <div class="sr-hero-title">🔬 AI-Powered Smart Reader</div>
        <div class="sr-hero-sub">
            Upload any <strong>PDF</strong>, <strong>Image</strong>, <strong>PowerPoint</strong>, or <strong>Video</strong>
            and unleash 40+ AI superpowers — summarize, quiz, flashcards, OCR, translate, Q&amp;A,
            diagram analysis, and much more. Instant. Intelligent. Incredible.
        </div>
    </div>
    """, unsafe_allow_html=True)

    # Format overview pills
    st.markdown("""
    <div class="sr-format-row">
        <div class="sr-format-pill">📄 PDF &nbsp;·&nbsp; 15 functions</div>
        <div class="sr-format-pill">🖼️ Image &nbsp;·&nbsp; 12 functions</div>
        <div class="sr-format-pill">📊 PowerPoint &nbsp;·&nbsp; 10 functions</div>
        <div class="sr-format-pill">🎬 Video &nbsp;·&nbsp; 6 functions</div>
    </div>
    """, unsafe_allow_html=True)

    # Format tabs
    tab_pdf, tab_img, tab_ppt, tab_vid = st.tabs([
        "📄 PDF Reader",
        "🖼️ Image AI",
        "📊 PPT Analyzer",
        "🎬 Video AI",
    ])

    with tab_pdf:
        _render_pdf_tab()

    with tab_img:
        _render_img_tab()

    with tab_ppt:
        _render_ppt_tab()

    with tab_vid:
        _render_video_tab()


# ── FREE API ADDITIONS ───────────────────────────────────────────────────────

def enrich_topic_with_wiki(topic: str) -> dict:
    """Fetch Wikipedia summary to enrich document context (free, no key)."""
    import urllib.request, urllib.parse, json
    try:
        url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{urllib.parse.quote(topic)}"
        req = urllib.request.Request(url, headers={"User-Agent": "ExamHelp/1.0"})
        with urllib.request.urlopen(req, timeout=5) as r:
            data = json.loads(r.read().decode())
        return {"title": data.get("title", ""), "summary": data.get("extract", "")[:500],
                "url": data.get("content_urls", {}).get("desktop", {}).get("page", "")}
    except Exception:
        return {}


def search_arxiv_papers(query: str, max_results: int = 5) -> list:
    """Search arXiv for related research papers (free, no key)."""
    import urllib.request, urllib.parse, re
    try:
        url = f"https://export.arxiv.org/api/query?search_query=all:{urllib.parse.quote(query)}&max_results={max_results}&sortBy=relevance"
        req = urllib.request.Request(url, headers={"User-Agent": "ExamHelp/1.0"})
        with urllib.request.urlopen(req, timeout=8) as r:
            text = r.read().decode()
        entries = re.findall(r"<entry>(.*?)</entry>", text, re.DOTALL)
        papers = []
        for entry in entries[:max_results]:
            title_m = re.search(r"<title>(.*?)</title>", entry)
            link_m = re.search(r"<id>(.*?)</id>", entry)
            summ_m = re.search(r"<summary>(.*?)</summary>", entry, re.DOTALL)
            if title_m and link_m:
                papers.append({
                    "title": title_m.group(1).strip().replace("\n", " "),
                    "url": link_m.group(1).strip(),
                    "abstract": summ_m.group(1).strip()[:250] if summ_m else "",
                })
        return papers
    except Exception:
        return []


def lookup_crossref_doi(query: str) -> list:
    """Search CrossRef for academic citations (free, no key)."""
    import urllib.request, urllib.parse, json
    try:
        url = f"https://api.crossref.org/works?query={urllib.parse.quote(query)}&rows=5"
        req = urllib.request.Request(url, headers={"User-Agent": "ExamHelp/1.0 (mailto:help@examhelp.ai)"})
        with urllib.request.urlopen(req, timeout=7) as r:
            data = json.loads(r.read().decode())
        items = data.get("message", {}).get("items", [])
        return [{"title": i.get("title", [""])[0], "doi": i.get("DOI", ""),
                 "url": f"https://doi.org/{i.get('DOI','')}", "year": i.get("published-print", {}).get("date-parts", [[""]])[0][0]}
                for i in items[:5] if i.get("title")]
    except Exception:
        return []


def get_reading_time(text: str, wpm: int = 200) -> dict:
    """Calculate estimated reading time for extracted document text."""
    words = len(text.split())
    minutes = max(1, round(words / wpm))
    return {"words": words, "minutes": minutes,
            "label": f"~{minutes} min read" if minutes < 60 else f"~{minutes // 60}h {minutes % 60}m read"}

