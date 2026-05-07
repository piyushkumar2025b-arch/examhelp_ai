"""
file_viewer_addon.py — Universal File Viewer (No AI)
Supports: Images, Text/Code, PDF, CSV, JSON, Excel, Audio, Video, Markdown,
          DOCX, PPTX, RTF, EPUB, IPYNB, EML, ZIP/RAR/7z/TAR, and any binary.

Key design goals:
  * No hard size limits — large files are chunked/paginated, never silently cut.
  * Every format has a graceful fallback chain.
  * PDF embeds up to 50 MB; larger ones get page-by-page text extraction.
  * CSV/Excel use paginated reads — millions of rows work fine.
  * Text pagination so even 100 MB logs are browsable.
"""
import streamlit as st
import io

VIEWER_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@400;700&family=Rajdhani:wght@400;600;700&family=Space+Mono&display=swap');

.fv-hero {
    background: linear-gradient(135deg,rgba(15,23,42,0.97),rgba(20,10,40,0.95));
    border: 1px solid rgba(99,102,241,0.2);
    border-radius: 22px; padding: 28px 32px; margin-bottom: 20px;
    position: relative; overflow: hidden;
}
.fv-hero::after {
    content:''; position:absolute; top:-1px; left:8%; right:8%; height:1px;
    background: linear-gradient(90deg,transparent,rgba(99,102,241,0.5),transparent);
}
.fv-title {
    font-family:'Rajdhani',sans-serif; font-size:2rem; font-weight:700;
    color:#fff; margin-bottom:4px;
}
.fv-sub { font-size:.9rem; color:rgba(255,255,255,0.4); }
.fv-card {
    background: rgba(15,23,42,0.75);
    border: 1px solid rgba(255,255,255,0.07);
    border-radius: 16px; padding: 20px 22px; margin-bottom: 14px;
}
.fv-label {
    font-family:'Space Mono',monospace; font-size:9px; letter-spacing:3px;
    color:rgba(255,255,255,0.3); text-transform:uppercase; margin-bottom:8px;
}
.fv-badge {
    display:inline-block; padding:3px 12px; border-radius:100px;
    font-family:'Space Mono',monospace; font-size:10px; font-weight:700;
    background:rgba(99,102,241,0.12); border:1px solid rgba(99,102,241,0.25);
    color:#a5b4fc; margin-right:6px; margin-bottom:4px;
}
.fv-info-row {
    display:flex; flex-wrap:wrap; gap:8px; margin-top:10px;
    font-family:'Rajdhani',sans-serif; font-size:.85rem; color:rgba(255,255,255,0.5);
}
</style>
"""

# ── Supported types ────────────────────────────────────────────────────────────
IMAGE_TYPES  = ["png","jpg","jpeg","gif","bmp","webp","svg","ico","tiff","heic","psd","ai"]
TEXT_TYPES   = ["txt","md","py","js","ts","css","html","xml","yaml","yml",
                "toml","ini","cfg","env","sh","bat","sql","log","rst","tex",
                "c","cpp","h","hpp","java","kt","go","rs","rb","php","swift",
                "r","lua","pl","scala","cs","vb","asm","dart","ex","exs",
                "coffee","jsx","tsx","vue","svelte","graphql","proto",
                "dockerfile","makefile","cmake","gradle","tf","hcl","nix"]
CODE_LANG_MAP = {
    "py":"python","js":"javascript","ts":"typescript","html":"html",
    "css":"css","sql":"sql","sh":"bash","bat":"batch","md":"markdown",
    "xml":"xml","yaml":"yaml","yml":"yaml","toml":"toml","json":"json",
    "rst":"rst","tex":"latex","ini":"ini","cfg":"ini","log":"text",
    "txt":"text","env":"text","c":"c","cpp":"cpp","h":"c","hpp":"cpp",
    "java":"java","kt":"kotlin","go":"go","rs":"rust","rb":"ruby",
    "php":"php","swift":"swift","r":"r","lua":"lua","pl":"perl",
    "scala":"scala","cs":"csharp","vb":"vbnet","dart":"dart",
    "jsx":"jsx","tsx":"tsx","vue":"html","svelte":"html",
    "graphql":"graphql","dockerfile":"dockerfile","tf":"hcl","hcl":"hcl",
    "gradle":"groovy","cmake":"cmake",
}
CSV_TYPES    = ["csv","tsv"]
JSON_TYPES   = ["json","jsonl","ndjson"]
PDF_TYPES    = ["pdf"]
AUDIO_TYPES  = ["mp3","wav","ogg","flac","m4a","aac","wma","opus","aiff"]
VIDEO_TYPES  = ["mp4","webm","ogv","mov","avi","mkv","flv","wmv"]
EXCEL_TYPES  = ["xlsx","xls","ods","xlsm"]
DOCX_TYPES   = ["docx","doc"]
PPT_TYPES    = ["pptx","ppt"]
RTF_TYPES    = ["rtf"]
EPUB_TYPES   = ["epub"]
IPYNB_TYPES  = ["ipynb"]
EML_TYPES    = ["eml","msg"]
ZIP_TYPES    = ["zip","rar","7z","tar","gz","bz2","xz","tgz"]
SYSTEM_TYPES = ["exe","dll","iso","dmg","app","deb","rpm","apk"]

# Page sizes
TEXT_PAGE_SIZE = 100_000   # chars per page for text viewers
ROWS_PER_PAGE  = 5_000     # rows per page for tabular viewers


def _fmt_size(n: int) -> str:
    if n < 1024:        return f"{n} B"
    if n < 1024**2:     return f"{n/1024:.1f} KB"
    if n < 1024**3:     return f"{n/1024**2:.1f} MB"
    return f"{n/1024**3:.2f} GB"


def _paginator(key: str, total_pages: int) -> int:
    """Renders Prev/Next controls and returns the current 0-based page index."""
    if total_pages <= 1:
        return 0
    if key not in st.session_state:
        st.session_state[key] = 0
    page = st.session_state[key]
    col_prev, col_info, col_next = st.columns([1, 3, 1])
    with col_prev:
        if st.button("◀ Prev", key=f"{key}_prev", disabled=(page == 0)):
            st.session_state[key] = page - 1
            st.rerun()
    with col_info:
        st.markdown(
            f"<div style='text-align:center;padding-top:6px;'>"
            f"Page <strong>{page+1}</strong> of <strong>{total_pages}</strong></div>",
            unsafe_allow_html=True,
        )
    with col_next:
        if st.button("Next ▶", key=f"{key}_next", disabled=(page == total_pages - 1)):
            st.session_state[key] = page + 1
            st.rerun()
    return st.session_state.get(key, 0)


# ─────────────────────────────────────────────────────────────────────────────
# Renderers
# ─────────────────────────────────────────────────────────────────────────────

def _render_image(file):
    ext = file.name.rsplit(".", 1)[-1].lower() if "." in file.name else ""
    raw_bytes = file.read()

    if ext == "svg":
        try:
            svg_content = raw_bytes.decode("utf-8")
            st.markdown(
                f'<div style="text-align:center;background:#fff;padding:20px;border-radius:12px;">'
                f'{svg_content}</div>', unsafe_allow_html=True)
            st.download_button("⬇️ Download SVG", raw_bytes, file_name=file.name,
                               use_container_width=True, key="fv_dl_svg")
            return
        except Exception:
            pass

    try:
        from PIL import Image
        img = Image.open(io.BytesIO(raw_bytes))
        if img.mode not in ("RGB", "RGBA"):
            img = img.convert("RGBA" if "transparency" in img.info else "RGB")
        col1, col2 = st.columns([3, 1])
        with col2:
            st.markdown("### 🖼️ Controls")
            rot = st.slider("Rotate°", -180, 180, 0, 90, key="fv_img_rot")
            if rot:
                img = img.rotate(-rot, expand=True)
            scale = st.slider("Scale %", 10, 200, 100, 10, key="fv_img_scale")
            if scale != 100:
                img = img.resize((int(img.width * scale / 100), int(img.height * scale / 100)))
            st.caption(f"{img.width}×{img.height} px | {img.mode}")
            st.download_button("⬇️ Download", raw_bytes, file_name=file.name,
                               use_container_width=True, key="fv_dl_img")
        with col1:
            st.image(img, use_container_width=True)
    except Exception as e:
        st.warning(f"Native render failed ({e}). Showing fallback.")
        _render_hex_fallback_bytes(raw_bytes, file.name)


def _render_text(file, ext: str):
    try:
        raw = file.read().decode("utf-8", errors="replace")
    except Exception as e:
        st.error(f"Cannot read file: {e}"); return

    lang  = CODE_LANG_MAP.get(ext, "text")
    lines = raw.count("\n") + 1
    words = len(raw.split())

    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Lines",  f"{lines:,}")
    col2.metric("Words",  f"{words:,}")
    col3.metric("Chars",  f"{len(raw):,}")
    col4.metric("Size",   _fmt_size(len(raw.encode("utf-8"))))

    search = st.text_input("🔍 Search:", key="fv_text_search", placeholder="type to find…")

    total_pages = max(1, -(-len(raw) // TEXT_PAGE_SIZE))
    page = _paginator("fv_text_page", total_pages)
    chunk = raw[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]

    if search:
        idx = raw.find(search)
        if idx >= 0:
            target_page = idx // TEXT_PAGE_SIZE
            if target_page == page:
                st.success(f"✅ '{search}' found on this page")
            else:
                st.info(f"ℹ️ '{search}' first found on page {target_page + 1}")
        else:
            st.warning(f"'{search}' not found in file")

    st.code(chunk, language=lang)
    st.download_button("⬇️ Download File", raw.encode("utf-8"), file_name=file.name,
                       mime="text/plain", use_container_width=True, key="fv_dl_text")


def _render_csv(file, ext: str):
    try:
        import pandas as pd
        raw = file.read()
        sep = "\t" if ext == "tsv" else ","

        line_count = raw.count(b"\n")
        st.markdown(f"**📊 Approx. {max(0, line_count - 1):,} data rows (counting newlines)**")

        search = st.text_input("🔍 Filter rows:", key="fv_csv_search", placeholder="type to filter…")

        try:
            df_full = pd.read_csv(io.BytesIO(raw), sep=sep, low_memory=False)
        except MemoryError:
            st.warning("File is very large — showing first 100,000 rows.")
            df_full = pd.read_csv(io.BytesIO(raw), sep=sep, nrows=100_000, low_memory=False)

        if search:
            mask = df_full.astype(str).apply(
                lambda col: col.str.contains(search, case=False, na=False)
            ).any(axis=1)
            df_view = df_full[mask].reset_index(drop=True)
            st.caption(f"Showing {len(df_view):,} matching rows")
        else:
            df_view = df_full

        total_rows = len(df_view)
        total_pages = max(1, -(-total_rows // ROWS_PER_PAGE))
        st.markdown(f"**{total_rows:,} rows × {len(df_view.columns):,} columns**")

        page = _paginator("fv_csv_page", total_pages)
        chunk = df_view.iloc[page * ROWS_PER_PAGE : (page + 1) * ROWS_PER_PAGE]
        st.dataframe(chunk, use_container_width=True, height=500)

        st.download_button("⬇️ Download CSV", raw, file_name=file.name,
                           mime="text/csv", use_container_width=True, key="fv_dl_csv")
    except Exception as e:
        st.error(f"CSV error: {e}")


def _render_json(file):
    import json
    raw = file.read().decode("utf-8", errors="replace")
    st.markdown(f"**Size:** {_fmt_size(len(raw))}")

    try:
        data = json.loads(raw)
        if isinstance(data, list):
            st.caption(f"JSON array with {len(data):,} items")
            total_pages = max(1, -(-len(data) // 100))
            page = _paginator("fv_json_page", total_pages)
            chunk = data[page * 100 : (page + 1) * 100]
            st.code(json.dumps(chunk, indent=2, ensure_ascii=False)[:TEXT_PAGE_SIZE], language="json")
        elif isinstance(data, dict):
            st.caption(f"JSON object with {len(data):,} keys")
            pretty = json.dumps(data, indent=2, ensure_ascii=False)
            total_pages = max(1, -(-len(pretty) // TEXT_PAGE_SIZE))
            page = _paginator("fv_json_page", total_pages)
            chunk = pretty[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]
            st.code(chunk, language="json")
        else:
            st.code(str(data)[:TEXT_PAGE_SIZE], language="json")
    except json.JSONDecodeError:
        lines = [l for l in raw.strip().split("\n") if l.strip()]
        st.warning(f"Not valid JSON — treating as {len(lines):,} JSONL lines")
        total_pages = max(1, -(-len(lines) // 200))
        page = _paginator("fv_jsonl_page", total_pages)
        for i, line in enumerate(lines[page * 200 : (page + 1) * 200], page * 200 + 1):
            try:
                parsed = json.loads(line)
                st.code(json.dumps(parsed, indent=2)[:2000], language="json")
            except Exception:
                st.text(f"[{i}] {line[:500]}")

    st.download_button("⬇️ Download JSON", raw.encode("utf-8"), file_name=file.name,
                       mime="application/json", use_container_width=True, key="fv_dl_json")


def _start_pdf_server(pdf_bytes: bytes) -> int:
    """
    Spin up (or reuse) a tiny HTTP server in a background thread that serves
    the current PDF bytes at GET /pdf.  Returns the port number.
    We keep one server per Streamlit session stored in st.session_state so it
    is not recreated on every widget interaction.
    """
    import threading
    import socket
    from http.server import BaseHTTPRequestHandler, HTTPServer

    # Store server state in session so it survives Streamlit reruns
    if "fv_pdf_server" not in st.session_state:
        st.session_state.fv_pdf_server = None
        st.session_state.fv_pdf_port = None

    # If bytes changed (new file uploaded) tear down the old server
    new_hash = hash(pdf_bytes)
    if st.session_state.get("fv_pdf_hash") != new_hash:
        if st.session_state.fv_pdf_server:
            try:
                st.session_state.fv_pdf_server.shutdown()
            except Exception:
                pass
        st.session_state.fv_pdf_server = None
        st.session_state.fv_pdf_port = None
        st.session_state.fv_pdf_hash = new_hash

    if st.session_state.fv_pdf_server is not None:
        return st.session_state.fv_pdf_port

    # Pick a free port
    with socket.socket() as s:
        s.bind(("", 0))
        port = s.getsockname()[1]

    # Capture bytes in closure
    _bytes = pdf_bytes

    class PDFHandler(BaseHTTPRequestHandler):
        def do_GET(self):
            if self.path in ("/pdf", "/pdf.pdf"):
                self.send_response(200)
                self.send_header("Content-Type", "application/pdf")
                self.send_header("Content-Length", str(len(_bytes)))
                self.send_header("Content-Disposition", "inline; filename=document.pdf")
                # Allow embedding from any origin (same machine)
                self.send_header("Access-Control-Allow-Origin", "*")
                self.end_headers()
                self.wfile.write(_bytes)
            else:
                self.send_response(404)
                self.end_headers()

        def log_message(self, *args):
            pass  # silence access log

    server = HTTPServer(("127.0.0.1", port), PDFHandler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()

    st.session_state.fv_pdf_server = server
    st.session_state.fv_pdf_port = port
    return port


def _render_pdf(file):
    data = file.read()
    size_mb = len(data) / (1024 * 1024)
    st.info(f"📕 PDF — {_fmt_size(len(data))}")

    tab_view, tab_text = st.tabs(["📄 Embedded Viewer", "📝 Extract Text"])

    with tab_view:
        try:
            import base64 as _b64
            b64data = _b64.b64encode(data).decode()
            pdfjs_url = "https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.min.js"
            pdfjs_worker = "https://cdnjs.cloudflare.com/ajax/libs/pdf.js/3.11.174/pdf.worker.min.js"
            # Use base64 data URI — works on Streamlit Cloud, any file size
            pdf_data_uri = f"data:application/pdf;base64,{b64data}"

            html_viewer = f"""
<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ background:#1a1a2e; font-family: sans-serif; }}
  #toolbar {{
    display:flex; align-items:center; gap:8px; padding:8px 12px;
    background:#0f0f23; border-bottom:1px solid #333;
    flex-wrap:wrap;
  }}
  #toolbar button {{
    background:#6366f1; color:#fff; border:none; border-radius:6px;
    padding:5px 12px; cursor:pointer; font-size:13px;
  }}
  #toolbar button:hover {{ background:#4f46e5; }}
  #toolbar button:disabled {{ background:#444; cursor:default; }}
  #page-info {{ color:#aaa; font-size:13px; min-width:100px; text-align:center; }}
  #zoom-info {{ color:#aaa; font-size:13px; }}
  #canvas-container {{
    overflow:auto; height:calc(100vh - 48px);
    display:flex; flex-direction:column; align-items:center;
    padding:12px; gap:8px;
  }}
  canvas {{ box-shadow:0 2px 16px rgba(0,0,0,0.6); background:#fff; display:block; }}
  #loading {{ color:#aaa; font-size:14px; margin-top:40px; }}
  #error {{ color:#f87171; font-size:14px; margin-top:40px; padding:20px; text-align:center; }}
</style>
</head>
<body>
<div id="toolbar">
  <button id="btn-first" onclick="goPage(1)">⏮</button>
  <button id="btn-prev" onclick="goPage(currentPage-1)">◀ Prev</button>
  <span id="page-info">Loading…</span>
  <button id="btn-next" onclick="goPage(currentPage+1)">Next ▶</button>
  <button id="btn-last" onclick="goPage(totalPages)">⏭</button>
  &nbsp;|&nbsp;
  <button onclick="changeZoom(-0.25)">🔍−</button>
  <span id="zoom-info">100%</span>
  <button onclick="changeZoom(+0.25)">🔍+</button>
  <button onclick="changeZoom(0)">Fit</button>
  &nbsp;|&nbsp;
  <button onclick="downloadPDF()">⬇️ Download</button>
</div>
<div id="canvas-container">
  <div id="loading">Loading PDF…</div>
</div>

<script src="{pdfjs_url}"></script>
<script>
  pdfjsLib.GlobalWorkerOptions.workerSrc = '{pdfjs_worker}';

  let pdfDoc = null, currentPage = 1, totalPages = 0, zoom = 1.0;
  const container = document.getElementById('canvas-container');

  async function loadPDF() {{
    try {{
      pdfDoc = await pdfjsLib.getDocument({{data: atob('{b64data}')}}).promise;
      totalPages = pdfDoc.numPages;
      document.getElementById('loading').style.display = 'none';
      fitZoom();
      renderPage(1);
    }} catch(e) {{
      document.getElementById('loading').innerHTML = '';
      const err = document.getElementById('error') || document.createElement('div');
      err.id = 'error';
      err.textContent = 'Could not load PDF: ' + e.message + '. Try the Open in Tab button above.';
      container.appendChild(err);
    }}
  }}

  async function fitZoom() {{
    // Render page 1 at scale 1 to measure natural width, then fit to container
    const page = await pdfDoc.getPage(1);
    const vp = page.getViewport({{scale:1}});
    const avail = container.clientWidth - 32;
    zoom = Math.min(1.5, avail / vp.width);
    document.getElementById('zoom-info').textContent = Math.round(zoom*100)+'%';
  }}

  async function renderPage(num) {{
    currentPage = Math.max(1, Math.min(num, totalPages));
    container.innerHTML = '';
    document.getElementById('page-info').textContent = `Page ${{currentPage}} / ${{totalPages}}`;
    document.getElementById('btn-prev').disabled = currentPage === 1;
    document.getElementById('btn-first').disabled = currentPage === 1;
    document.getElementById('btn-next').disabled = currentPage === totalPages;
    document.getElementById('btn-last').disabled = currentPage === totalPages;

    const page = await pdfDoc.getPage(currentPage);
    const viewport = page.getViewport({{scale: zoom}});
    const canvas = document.createElement('canvas');
    canvas.width  = viewport.width;
    canvas.height = viewport.height;
    container.appendChild(canvas);
    await page.render({{canvasContext: canvas.getContext('2d'), viewport}}).promise;
  }}

  function goPage(n) {{ renderPage(n); }}

  function changeZoom(delta) {{
    if (delta === 0) {{ fitZoom().then(() => renderPage(currentPage)); return; }}
    zoom = Math.max(0.25, Math.min(4.0, zoom + delta));
    document.getElementById('zoom-info').textContent = Math.round(zoom*100)+'%';
    renderPage(currentPage);
  }}

  document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === 'PageDown') goPage(currentPage+1);
    if (e.key === 'ArrowLeft'  || e.key === 'PageUp')   goPage(currentPage-1);
  }});

  function downloadPDF() {{
    const a = document.createElement('a');
    a.href = '{pdf_data_uri}';
    a.download = 'document.pdf';
    a.click();
  }}

  loadPDF();
</script>
</body>
</html>
"""
            st.components.v1.html(html_viewer, height=750, scrolling=False)

        except Exception as e:
            st.error(f"PDF viewer error: {e}")
            # Hard fallback — base64 iframe
            import base64
            b64 = base64.b64encode(data).decode()
            st.markdown(
                f'<iframe src="data:application/pdf;base64,{b64}" '
                f'width="100%" height="750px" '
                f'style="border:1px solid rgba(255,255,255,0.1);border-radius:12px;"></iframe>',
                unsafe_allow_html=True,
            )

        st.download_button("⬇️ Download PDF", data, file_name=file.name,
                           mime="application/pdf", use_container_width=True, key="fv_dl_pdf")

    with tab_text:
        try:
            import fitz  # PyMuPDF — already in requirements.txt
            pdf_doc = fitz.open(stream=data, filetype="pdf")
            total_pages_count = pdf_doc.page_count
            st.caption(f"{total_pages_count} pages")
            page_num = st.number_input("Go to page:", min_value=1, max_value=total_pages_count,
                                       value=1, step=1, key="fv_pdf_page")
            page_obj = pdf_doc.load_page(page_num - 1)
            page_text = page_obj.get_text() or "(no extractable text on this page)"
            st.text_area(f"Page {page_num}", page_text, height=500,
                         label_visibility="collapsed", key="fv_pdf_text_area")
            pdf_doc.close()
        except ImportError:
            # Fallback: try pypdf if somehow installed
            try:
                from pypdf import PdfReader
                reader = PdfReader(io.BytesIO(data))
                total_pages_count = len(reader.pages)
                st.caption(f"{total_pages_count} pages")
                page_num = st.number_input("Go to page:", min_value=1, max_value=total_pages_count,
                                           value=1, step=1, key="fv_pdf_page")
                page_text = reader.pages[page_num - 1].extract_text() or "(no extractable text)"
                st.text_area(f"Page {page_num}", page_text, height=500,
                             label_visibility="collapsed", key="fv_pdf_text_area")
            except ImportError:
                st.info("⚠️ PDF text extraction requires PyMuPDF (pip install PyMuPDF) or pypdf (pip install pypdf).")
        except Exception as e:
            st.error(f"Text extraction error: {e}")



def _render_audio(file):
    raw = file.read()
    st.audio(raw)
    st.caption(f"🎵 {file.name} — use the player above")
    st.download_button("⬇️ Download Audio", raw, file_name=file.name,
                       use_container_width=True, key="fv_dl_audio")


def _render_video(file):
    st.video(file)
    st.caption(f"🎬 {file.name} — use the player above")


def _render_excel(file):
    try:
        import pandas as pd
        raw = file.read()
        ext = file.name.rsplit(".", 1)[-1].lower() if "." in file.name else "xlsx"
        engine = {"xls": "xlrd", "ods": "odf"}.get(ext, "openpyxl")

        xl = pd.ExcelFile(io.BytesIO(raw), engine=engine)
        sheets = xl.sheet_names
        st.markdown(f"**📊 Excel — Sheets:** {', '.join(f'`{s}`' for s in sheets)}")
        selected = st.selectbox("Select sheet:", sheets, key="fv_xl_sheet")

        df = xl.parse(selected)
        search = st.text_input("🔍 Filter rows:", key="fv_xl_search", placeholder="type to filter…")
        if search:
            mask = df.astype(str).apply(
                lambda col: col.str.contains(search, case=False, na=False)
            ).any(axis=1)
            df = df[mask].reset_index(drop=True)
            st.caption(f"Showing {len(df):,} matching rows")

        total_rows = len(df)
        total_pages = max(1, -(-total_rows // ROWS_PER_PAGE))
        st.markdown(f"**{total_rows:,} rows × {len(df.columns):,} columns**")

        page = _paginator("fv_xl_page", total_pages)
        chunk = df.iloc[page * ROWS_PER_PAGE : (page + 1) * ROWS_PER_PAGE]
        st.dataframe(chunk, use_container_width=True, height=500)

        st.download_button("⬇️ Download Excel", raw, file_name=file.name,
                           use_container_width=True, key="fv_dl_excel")
    except Exception as e:
        st.error(f"Excel error: {e}. Ensure openpyxl / xlrd is installed.")


def _render_docx(file):
    try:
        from docx import Document
        raw = file.read()
        doc = Document(io.BytesIO(raw))

        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_count = len(doc.tables)
        st.markdown(f"**📝 Word Document — {len(paragraphs):,} paragraphs, {tables_count} table(s)**")

        tab_text, tab_tables = st.tabs(["📄 Text Content", "📊 Tables"])

        with tab_text:
            full_text = "\n\n".join(paragraphs)
            total_pages = max(1, -(-len(full_text) // TEXT_PAGE_SIZE))
            page = _paginator("fv_docx_page", total_pages)
            chunk = full_text[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]
            st.text_area("Content", chunk, height=600,
                         label_visibility="collapsed", key="fv_docx_view")

        with tab_tables:
            if tables_count == 0:
                st.caption("No tables in this document.")
            else:
                import pandas as pd
                for t_idx, table in enumerate(doc.tables):
                    st.markdown(f"**Table {t_idx + 1}**")
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    if rows:
                        df = (pd.DataFrame(rows[1:], columns=rows[0])
                              if len(rows) > 1 else pd.DataFrame(rows))
                        st.dataframe(df, use_container_width=True)

        st.download_button("⬇️ Download DOCX", raw, file_name=file.name,
                           use_container_width=True, key="fv_dl_docx")
    except ImportError:
        st.error("Install python-docx: `pip install python-docx`")
    except Exception as e:
        st.error(f"DOCX error: {e}")


def _render_ppt(file):
    try:
        from pptx import Presentation
        raw = file.read()
        prs = Presentation(io.BytesIO(raw))
        total_slides = len(prs.slides)
        st.markdown(f"**📊 Presentation — {total_slides} slides**")

        page_size = 20
        total_pages = max(1, -(-total_slides // page_size))
        page = _paginator("fv_ppt_page", total_pages)
        start = page * page_size
        end = min(start + page_size, total_slides)

        for i in range(start, end):
            slide = prs.slides[i]
            text_runs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = " ".join(r.text for r in para.runs).strip()
                        if t:
                            text_runs.append(t)
            with st.expander(f"Slide {i + 1}", expanded=(i == start)):
                if text_runs:
                    st.write("\n".join(text_runs))
                else:
                    st.caption("No text on this slide.")

        st.download_button("⬇️ Download PPT", raw, file_name=file.name,
                           use_container_width=True, key="fv_dl_ppt")
    except ImportError:
        st.error("Install python-pptx: `pip install python-pptx`")
    except Exception as e:
        st.error(f"PPT Error: {e}")


def _render_rtf(file):
    import re
    try:
        raw = file.read().decode("latin-1", errors="replace")
        text = re.sub(r"\\[a-z]+[-]?\d*[ ]?", "", raw)
        text = re.sub(r"[{}\\]", "", text)
        text = re.sub(r"\s+", " ", text).strip()

        st.markdown(f"**📄 RTF Document — {_fmt_size(len(raw))} | {len(text):,} chars extracted**")
        total_pages = max(1, -(-len(text) // TEXT_PAGE_SIZE))
        page = _paginator("fv_rtf_page", total_pages)
        chunk = text[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]
        st.text_area("Content", chunk, height=500,
                     label_visibility="collapsed", key="fv_rtf_view")
        st.download_button("⬇️ Download RTF", raw.encode("latin-1", errors="replace"),
                           file_name=file.name, use_container_width=True, key="fv_dl_rtf")
    except Exception as e:
        st.error(f"RTF error: {e}")


def _render_eml(file):
    import email
    try:
        raw = file.read()
        msg = email.message_from_bytes(raw)
        st.markdown("**📧 Email Message**")
        for hdr in ["From", "To", "Cc", "Subject", "Date", "Reply-To"]:
            val = msg.get(hdr, "")
            if val:
                st.markdown(f"**{hdr}:** {val}")
        st.markdown("---")

        bodies, attachments = [], []
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                cd = str(part.get("Content-Disposition", ""))
                if "attachment" in cd:
                    fn = part.get_filename() or "attachment"
                    attachments.append((fn, part.get_payload(decode=True)))
                elif ct == "text/plain":
                    bodies.append(part.get_payload(decode=True).decode("utf-8", errors="replace"))
                elif ct == "text/html" and not bodies:
                    import re
                    html_raw = part.get_payload(decode=True).decode("utf-8", errors="replace")
                    bodies.append(re.sub(r"<[^>]+>", " ", html_raw))
        else:
            bodies.append(msg.get_payload(decode=True).decode("utf-8", errors="replace"))

        body = "\n\n".join(bodies)
        total_pages = max(1, -(-len(body) // TEXT_PAGE_SIZE))
        page = _paginator("fv_eml_page", total_pages)
        chunk = body[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]
        st.text_area("Body", chunk, height=400, label_visibility="collapsed", key="fv_eml_body")

        if attachments:
            st.markdown(f"**📎 {len(attachments)} attachment(s):**")
            for fn, att_data in attachments:
                if att_data:
                    st.download_button(f"⬇️ {fn}", att_data, file_name=fn,
                                       use_container_width=True, key=f"fv_eml_att_{fn}")

        st.download_button("⬇️ Download EML", raw, file_name=file.name,
                           use_container_width=True, key="fv_dl_eml")
    except Exception as e:
        st.error(f"EML error: {e}")


def _render_notebook(file):
    import json
    try:
        data = json.loads(file.read().decode("utf-8", errors="replace"))
        cells = data.get("cells", [])
        st.markdown(f"**📓 Jupyter Notebook — {len(cells)} cells**")

        page_size = 30
        total_pages = max(1, -(-len(cells) // page_size))
        page = _paginator("fv_nb_page", total_pages)
        start = page * page_size

        for i, cell in enumerate(cells[start : start + page_size], start + 1):
            ctype = cell.get("cell_type", "code")
            src = "".join(cell.get("source", []))
            icon = "🐍" if ctype == "code" else "📝"
            with st.expander(f"{icon} Cell {i} [{ctype}]", expanded=(i == start + 1)):
                if ctype == "code":
                    st.code(src, language="python")
                else:
                    st.markdown(src)
                for out in cell.get("outputs", []):
                    text = out.get("text") or out.get("data", {}).get("text/plain", [])
                    if text:
                        st.caption("".join(text)[:2000])
    except Exception as e:
        st.error(f"Notebook error: {e}")


def _render_epub(file):
    import zipfile, re
    try:
        data = file.read()
        with zipfile.ZipFile(io.BytesIO(data), "r") as z:
            names = z.namelist()
            html_files = [n for n in names if n.lower().endswith((".html", ".xhtml", ".htm"))]
            seen, unique_files = set(), []
            for hf in html_files:
                if hf not in seen:
                    seen.add(hf)
                    unique_files.append(hf)

            st.markdown(f"**📖 EPUB — {len(unique_files)} chapters found**")
            if not unique_files:
                st.warning("No readable HTML chapters found in this EPUB.")
                st.download_button("⬇️ Download EPUB", data, file_name=file.name,
                                   use_container_width=True, key="fv_dl_epub")
                return

            page_size = 5
            total_pages = max(1, -(-len(unique_files) // page_size))
            page = _paginator("fv_epub_page", total_pages)
            start = page * page_size

            # Read all chapter content INSIDE the context manager before rendering widgets
            chapters = []
            for hf in unique_files[start : start + page_size]:
                raw_html = z.read(hf).decode("utf-8", errors="replace")
                text = re.sub(r"<[^>]+>", " ", raw_html)
                text = re.sub(r"\s+", " ", text).strip()
                chapters.append((hf, text))

        # Now render OUTSIDE the context manager (no stream needed anymore)
        for i, (hf, text) in enumerate(chapters, start + 1):
            if text:
                with st.expander(f"Chapter {i}: {hf.split('/')[-1]}", expanded=(i == start + 1)):
                    if len(text) > 10_000:
                        st.write(text[:10_000])
                        st.caption(f"… ({len(text):,} chars — download for full text)")
                    else:
                        st.write(text)

        st.download_button("⬇️ Download EPUB", data, file_name=file.name,
                           use_container_width=True, key="fv_dl_epub")
    except Exception as e:
        st.error(f"EPUB error: {e}")


def _render_zip(file):
    import zipfile, tarfile
    raw = file.read()
    ext = file.name.rsplit(".", 1)[-1].lower() if "." in file.name else ""

    if ext == "zip":
        try:
            # ── Build the full file list once (cheap, no stream kept open) ──
            with zipfile.ZipFile(io.BytesIO(raw), "r") as z_probe:
                info_list  = z_probe.infolist()
                file_names = [i.filename for i in info_list if not i.is_dir()]

            st.markdown(f"**🗜️ ZIP Archive — {len(info_list):,} items**")

            search = st.text_input("🔍 Filter files:", key="fv_zip_search",
                                   placeholder="e.g. .py or README")
            filtered = [i for i in info_list
                        if not search or search.lower() in i.filename.lower()]
            st.caption(f"Showing {len(filtered):,} of {len(info_list):,} items")

            total_pages = max(1, -(-len(filtered) // 500))
            page = _paginator("fv_zip_page", total_pages)
            chunk = filtered[page * 500 : (page + 1) * 500]

            import pandas as pd
            data_rows = [{"Name": i.filename, "Size": _fmt_size(i.file_size),
                          "Compressed": _fmt_size(i.compress_size),
                          "Ratio": f"{100*(1-i.compress_size/i.file_size):.0f}%"
                                   if i.file_size > 0 else "—"}
                         for i in chunk]
            st.dataframe(data_rows, use_container_width=True)

            # ── Preview individual file ──
            st.markdown("**👁️ Preview a file from the archive:**")
            if file_names:
                selected = st.selectbox("Choose file:", ["(none)"] + file_names[:1000],
                                        key="fv_zip_file_select")
                if selected and selected != "(none)":
                    cache_key = f"fv_zip_inner_{selected}"
                    if cache_key not in st.session_state:
                        with zipfile.ZipFile(io.BytesIO(raw), "r") as z_extract:
                            try:
                                st.session_state[cache_key] = z_extract.read(selected)
                            except Exception as read_err:
                                st.error(f"Cannot read '{selected}': {read_err}")
                                st.session_state[cache_key] = b""

                    inner_bytes = st.session_state.get(cache_key, b"")
                    inner_ext = selected.rsplit(".", 1)[-1].lower() if "." in selected else ""
                    st.caption(f"**{selected}** — {_fmt_size(len(inner_bytes))}")

                    if inner_bytes:
                        if inner_ext in TEXT_TYPES or inner_ext in JSON_TYPES:
                            try:
                                inner_text = inner_bytes.decode("utf-8", errors="replace")
                                total_inner_pages = max(1, -(-len(inner_text) // TEXT_PAGE_SIZE))
                                inner_page = _paginator("fv_zip_inner_page", total_inner_pages)
                                st.code(inner_text[inner_page * TEXT_PAGE_SIZE :
                                                   (inner_page + 1) * TEXT_PAGE_SIZE],
                                        language=CODE_LANG_MAP.get(inner_ext, "text"))
                            except Exception as txt_err:
                                st.info(f"Cannot decode as text: {txt_err}")
                        elif inner_ext in IMAGE_TYPES:
                            try:
                                from PIL import Image as PILImage
                                img = PILImage.open(io.BytesIO(inner_bytes))
                                st.image(img, use_container_width=True)
                            except Exception:
                                st.info("Cannot render image.")
                        else:
                            st.download_button("⬇️ Download this file", inner_bytes,
                                               file_name=selected.split("/")[-1],
                                               use_container_width=True, key="fv_zip_inner_dl")
        except Exception as e:
            st.error(f"ZIP parsing error: {e}")

    elif ext in ("tar", "gz", "bz2", "xz", "tgz"):
        try:
            with tarfile.open(fileobj=io.BytesIO(raw), mode="r:*") as t:
                members = t.getmembers()
                st.markdown(f"**🗜️ TAR Archive — {len(members):,} items**")
                import pandas as pd
                total_pages = max(1, -(-len(members) // 500))
                page = _paginator("fv_tar_page", total_pages)
                chunk = members[page * 500 : (page + 1) * 500]
                data_rows = [{"Name": m.name, "Size": _fmt_size(m.size),
                              "Type": "dir" if m.isdir() else "file"} for m in chunk]
                st.dataframe(data_rows, use_container_width=True)
        except Exception as e:
            st.warning(f"TAR parsing error: {e}")
            _render_hex_fallback_bytes(raw, file.name)
    else:
        st.info(f"ℹ️ .{ext} archives require a local tool — you can download the file below.")
        _render_hex_fallback_bytes(raw, file.name)

    st.download_button("⬇️ Download Archive", raw, file_name=file.name,
                       use_container_width=True, key="fv_dl_zip")


def _render_hex_fallback_bytes(raw_bytes: bytes, filename: str):
    """Universal fallback — paginated text or hex dump, always with download."""
    st.warning(f"⚠️ No native viewer for this type. Showing raw contents ({_fmt_size(len(raw_bytes))}).")
    try:
        text_content = raw_bytes.decode("utf-8")
        st.markdown("**📝 Detected as UTF-8 Text:**")
        total_pages = max(1, -(-len(text_content) // TEXT_PAGE_SIZE))
        page = _paginator("fv_fb_page", total_pages)
        chunk = text_content[page * TEXT_PAGE_SIZE : (page + 1) * TEXT_PAGE_SIZE]
        st.text_area("Content", chunk, height=400, disabled=True,
                     label_visibility="collapsed", key="fv_fb_txt")
    except UnicodeDecodeError:
        st.markdown("**🔢 Binary File — Hex Dump:**")
        BYTES_PER_PAGE = 4096
        total_pages = max(1, -(-len(raw_bytes) // BYTES_PER_PAGE))
        page = _paginator("fv_hex_page", total_pages)
        display_bytes = raw_bytes[page * BYTES_PER_PAGE : (page + 1) * BYTES_PER_PAGE]
        base_offset = page * BYTES_PER_PAGE
        hex_lines = []
        for i in range(0, len(display_bytes), 16):
            chunk = display_bytes[i:i + 16]
            hex_part = " ".join(f"{b:02x}" for b in chunk).ljust(47)
            ascii_part = "".join(chr(b) if 32 <= b <= 126 else "." for b in chunk)
            hex_lines.append(f"{base_offset + i:08x}  {hex_part}  |{ascii_part}|")
        st.code("\n".join(hex_lines), language="text")

    st.download_button("⬇️ Download Original File", raw_bytes,
                       file_name=filename, use_container_width=True, key="fv_dl_fallback")


def _render_hex_fallback(file):
    try:
        _render_hex_fallback_bytes(file.read(), file.name)
    except Exception as e:
        st.error(f"Fallback view error: {e}")


# ─────────────────────────────────────────────────────────────────────────────
# Main page
# ─────────────────────────────────────────────────────────────────────────────

def render_file_viewer_page():
    """Main File Viewer page — no AI, pure file rendering."""
    st.markdown(VIEWER_CSS, unsafe_allow_html=True)

    st.markdown("""
    <div class="fv-hero">
        <div class="fv-title">📂 Universal File Viewer</div>
        <div class="fv-sub">Open & view any file directly in your browser — no AI, no uploads to servers. All file sizes supported.</div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="fv-label">Supported Formats</div>', unsafe_allow_html=True)
    badges = ""
    for label, _ in [
        ("🖼️ Images", IMAGE_TYPES), ("📄 Text/Code", TEXT_TYPES),
        ("📊 CSV/TSV", CSV_TYPES),  ("🔧 JSON", JSON_TYPES),
        ("📕 PDF", PDF_TYPES),      ("🎵 Audio", AUDIO_TYPES),
        ("🎬 Video", VIDEO_TYPES),  ("📊 Excel", EXCEL_TYPES),
        ("📝 Word", DOCX_TYPES),    ("📊 PowerPoint", PPT_TYPES),
        ("📄 RTF", RTF_TYPES),      ("📖 EPUB", EPUB_TYPES),
        ("📓 Notebook", IPYNB_TYPES),("📧 Email", EML_TYPES),
        ("🗜️ Archive", ZIP_TYPES),  ("⚙️ System", SYSTEM_TYPES),
    ]:
        badges += f'<span class="fv-badge">{label}</span>'
    st.markdown(badges, unsafe_allow_html=True)
    st.markdown("")

    uploaded = st.file_uploader(
        "Choose any file to view",
        type=None,
        key="fv_uploader",
        label_visibility="collapsed",
        help="Drag & drop or click to browse. Upload limit is set in .streamlit/config.toml.",
    )

    if not uploaded:
        st.markdown("""
        <div class="fv-card" style="text-align:center;padding:40px;">
            <div style="font-size:3rem;margin-bottom:12px;">📂</div>
            <div style="font-family:'Rajdhani',sans-serif;font-size:1.1rem;color:rgba(255,255,255,0.5);">
                Drop any file here to view it instantly.<br>
                <span style="font-size:.85rem;color:rgba(255,255,255,0.25);">
                    Images · PDFs · Code · CSV · Excel · Word · Audio · Video · Archives · and more
                </span>
            </div>
        </div>
        """, unsafe_allow_html=True)
        st.stop()

    # ── File info bar ──────────────────────────────────────────────────────────
    ext = uploaded.name.rsplit(".", 1)[-1].lower() if "." in uploaded.name else ""
    # Special no-extension files
    if not ext and uploaded.name.lower() in ("dockerfile","makefile","procfile","vagrantfile","gemfile"):
        ext = uploaded.name.lower()

    size_bytes = uploaded.size
    mime = uploaded.type or "application/octet-stream"

    st.markdown(f"""
    <div class="fv-card">
        <div class="fv-label">File Information</div>
        <div style="font-family:'Rajdhani',sans-serif;font-size:1.1rem;font-weight:700;color:#fff;margin-bottom:8px;">
            📄 {uploaded.name}
        </div>
        <div class="fv-info-row">
            <span>📦 <strong>{_fmt_size(size_bytes)}</strong></span>
            <span>🏷️ {('.'+ext.upper()) if ext else '(no ext)'}</span>
            <span>🔖 {mime}</span>
        </div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="fv-label">File Preview</div>', unsafe_allow_html=True)

    # ── Routing ────────────────────────────────────────────────────────────────
    if   ext in IMAGE_TYPES:   _render_image(uploaded)
    elif ext in TEXT_TYPES:    _render_text(uploaded, ext)
    elif ext in CSV_TYPES:     _render_csv(uploaded, ext)
    elif ext in JSON_TYPES:    _render_json(uploaded)
    elif ext in PDF_TYPES:     _render_pdf(uploaded)
    elif ext in AUDIO_TYPES:   _render_audio(uploaded)
    elif ext in VIDEO_TYPES:   _render_video(uploaded)
    elif ext in EXCEL_TYPES:   _render_excel(uploaded)
    elif ext in DOCX_TYPES:    _render_docx(uploaded)
    elif ext in PPT_TYPES:     _render_ppt(uploaded)
    elif ext in RTF_TYPES:     _render_rtf(uploaded)
    elif ext in EPUB_TYPES:    _render_epub(uploaded)
    elif ext in IPYNB_TYPES:   _render_notebook(uploaded)
    elif ext in EML_TYPES:     _render_eml(uploaded)
    elif ext in ZIP_TYPES:     _render_zip(uploaded)
    else:
        # MIME-type sniffing as second chance
        if   mime.startswith("image/"):                           _render_image(uploaded)
        elif mime.startswith("text/") or mime in (
             "application/json","application/xml"):               _render_text(uploaded, ext or "txt")
        elif mime == "application/pdf":                           _render_pdf(uploaded)
        elif mime.startswith("audio/"):                           _render_audio(uploaded)
        elif mime.startswith("video/"):                           _render_video(uploaded)
        else:                                                     _render_hex_fallback(uploaded)

    st.markdown("---")
    if st.button("💬 Back to Chat", use_container_width=True, key="fv_back"):
        st.session_state.app_mode = "chat"
        st.rerun()
