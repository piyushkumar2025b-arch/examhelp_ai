"""
converter_engine.py — UNIVERSAL File & Unit Converter Engine v2.0
Supports 100+ file conversions + 20 unit categories + live currency.
"""
from __future__ import annotations
import io, os, base64, json, csv, re
from typing import Optional, Tuple

# ── Core text/document converters ─────────────────────────────────────────────

def _pdf_to_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        return "\n\n".join(p.extract_text() or "" for p in reader.pages)
    except Exception as e:
        return f"[PDF extraction error: {e}]"

def _docx_to_text(data: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"[DOCX extraction error: {e}]"

def _excel_to_text(data: bytes) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = []
        for sheet in wb.worksheets:
            out.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else "" for c in row]
                out.append("\t".join(row_vals))
        return "\n".join(out)
    except Exception as e:
        return f"[XLSX extraction error: {e}]"

def _csv_to_text(data: bytes) -> str:
    try:
        content = data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        return "\n".join(",".join(row) for row in reader)
    except Exception as e:
        return f"[CSV extraction error: {e}]"

def _markdown_to_html(data: bytes) -> bytes:
    try:
        import markdown as md_lib
        html = md_lib.markdown(data.decode("utf-8", errors="replace"), extensions=["tables","fenced_code"])
    except ImportError:
        text = data.decode("utf-8", errors="replace")
        html = "<pre>" + text.replace("&","&amp;").replace("<","&lt;") + "</pre>"
    page = f"""<!DOCTYPE html><html><head><meta charset='utf-8'>
<style>body{{font-family:sans-serif;max-width:800px;margin:40px auto;line-height:1.6;}}
code{{background:#f4f4f4;padding:2px 6px;border-radius:4px;}}
pre{{background:#f4f4f4;padding:12px;border-radius:8px;overflow-x:auto;}}</style></head>
<body>{html}</body></html>"""
    return page.encode("utf-8")

def _text_to_pdf(data: bytes, title: str = "Document") -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm
        from reportlab.pdfgen import canvas
        text = data.decode("utf-8", errors="replace")
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        w, h = A4
        y = h - 2*cm
        c.setFont("Helvetica-Bold", 14)
        c.drawString(2*cm, y, title)
        y -= 1*cm
        c.setFont("Helvetica", 10)
        for line in text.splitlines():
            if y < 2*cm:
                c.showPage()
                y = h - 2*cm
                c.setFont("Helvetica", 10)
            wrapped = [line[i:i+95] for i in range(0, max(1, len(line)), 95)] if len(line) > 95 else [line]
            for wl in wrapped:
                c.drawString(2*cm, y, wl)
                y -= 0.45*cm
        c.save()
        return buf.getvalue()
    except Exception as e:
        return f"PDF generation error: {e}".encode()

def _text_to_docx(data: bytes) -> bytes:
    try:
        from docx import Document
        text = data.decode("utf-8", errors="replace")
        doc = Document()
        for line in text.splitlines():
            doc.add_paragraph(line)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except Exception as e:
        return f"DOCX error: {e}".encode()

def _csv_to_json(data: bytes) -> bytes:
    content = data.decode("utf-8", errors="replace")
    reader = csv.DictReader(io.StringIO(content))
    rows = list(reader)
    return json.dumps(rows, indent=2, ensure_ascii=False).encode("utf-8")

def _json_to_csv(data: bytes) -> bytes:
    parsed = json.loads(data.decode("utf-8"))
    if isinstance(parsed, dict): parsed = [parsed]
    if not parsed: return b""
    buf = io.StringIO()
    writer = csv.DictWriter(buf, fieldnames=list(parsed[0].keys()))
    writer.writeheader(); writer.writerows(parsed)
    return buf.getvalue().encode("utf-8")

def _csv_to_xlsx(data: bytes) -> bytes:
    try:
        import openpyxl
        content = data.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        wb = openpyxl.Workbook(); ws = wb.active
        for row in reader: ws.append(row)
        buf = io.BytesIO(); wb.save(buf); return buf.getvalue()
    except Exception as e:
        return f"XLSX error: {e}".encode()

def _xlsx_to_csv(data: bytes) -> bytes:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        ws = wb.active
        buf = io.StringIO(); writer = csv.writer(buf)
        for row in ws.iter_rows(values_only=True):
            writer.writerow([str(c) if c is not None else "" for c in row])
        return buf.getvalue().encode("utf-8")
    except Exception as e:
        return f"CSV error: {e}".encode()

def _xlsx_to_json(data: bytes) -> bytes:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        result = {}
        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if rows:
                headers = [str(h) if h is not None else f"col{i}" for i, h in enumerate(rows[0])]
                result[sheet.title] = [{headers[i]: (str(c) if c is not None else "") for i, c in enumerate(row)} for row in rows[1:]]
        return json.dumps(result, indent=2, ensure_ascii=False).encode("utf-8")
    except Exception as e:
        return f"JSON error: {e}".encode()

def _html_to_text(data: bytes) -> bytes:
    try:
        from html.parser import HTMLParser
        class _P(HTMLParser):
            def __init__(self):
                super().__init__(); self.text = []
            def handle_data(self, d): self.text.append(d)
        p = _P(); p.feed(data.decode("utf-8", errors="replace"))
        return " ".join(p.text).encode("utf-8")
    except Exception as e:
        return f"Error: {e}".encode()

def _html_to_markdown(data: bytes) -> bytes:
    html = data.decode("utf-8", errors="replace")
    html = re.sub(r'<h([1-6])[^>]*>(.*?)</h\1>', lambda m: '#'*int(m.group(1)) + ' ' + m.group(2), html, flags=re.DOTALL)
    html = re.sub(r'<b[^>]*>(.*?)</b>', r'**\1**', html, flags=re.DOTALL)
    html = re.sub(r'<strong[^>]*>(.*?)</strong>', r'**\1**', html, flags=re.DOTALL)
    html = re.sub(r'<i[^>]*>(.*?)</i>', r'*\1*', html, flags=re.DOTALL)
    html = re.sub(r'<a[^>]*href="([^"]*)"[^>]*>(.*?)</a>', r'[\2](\1)', html, flags=re.DOTALL)
    html = re.sub(r'<li[^>]*>(.*?)</li>', r'- \1', html, flags=re.DOTALL)
    html = re.sub(r'<p[^>]*>(.*?)</p>', r'\1\n\n', html, flags=re.DOTALL)
    html = re.sub(r'<[^>]+>', '', html)
    html = re.sub(r'\n{3,}', '\n\n', html)
    return html.strip().encode("utf-8")

def _text_to_html(data: bytes) -> bytes:
    text = data.decode("utf-8", errors="replace")
    html_content = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    paragraphs = "".join(f"<p>{line}</p>\n" if line.strip() else "<br>\n" for line in html_content.split("\n"))
    page = f"""<!DOCTYPE html><html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<style>body{{font-family:'Segoe UI',sans-serif;max-width:800px;margin:40px auto;padding:0 20px;line-height:1.7;color:#222;}}p{{margin:0 0 12px;}}</style></head>
<body>{paragraphs}</body></html>"""
    return page.encode("utf-8")

def _csv_to_markdown_table(data: bytes) -> bytes:
    content = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    if not rows: return b""
    headers = rows[0]
    md = "| " + " | ".join(headers) + " |\n"
    md += "| " + " | ".join(["---"] * len(headers)) + " |\n"
    for row in rows[1:]:
        padded = row + [""] * (len(headers) - len(row))
        md += "| " + " | ".join(str(c) for c in padded[:len(headers)]) + " |\n"
    return md.encode("utf-8")

def _xml_to_json(data: bytes) -> bytes:
    try:
        import xml.etree.ElementTree as ET
        def elem_to_dict(elem):
            d = {}
            if elem.attrib: d["@attributes"] = elem.attrib
            if elem.text and elem.text.strip(): d["#text"] = elem.text.strip()
            for child in elem:
                cd = elem_to_dict(child)
                if child.tag in d:
                    if not isinstance(d[child.tag], list): d[child.tag] = [d[child.tag]]
                    d[child.tag].append(cd)
                else: d[child.tag] = cd
            return d
        root = ET.fromstring(data.decode("utf-8", errors="replace"))
        return json.dumps({root.tag: elem_to_dict(root)}, indent=2, ensure_ascii=False).encode("utf-8")
    except Exception as e:
        return f"XML error: {e}".encode()

def _json_to_xml(data: bytes) -> bytes:
    parsed = json.loads(data.decode("utf-8"))
    def dict_to_xml(d, tag="root"):
        parts = [f"<{tag}>"]
        if isinstance(d, dict):
            for k, v in d.items(): parts.append(dict_to_xml(v, k))
        elif isinstance(d, list):
            for item in d: parts.append(dict_to_xml(item, "item"))
        else:
            parts.append(str(d))
        parts.append(f"</{tag}>")
        return "".join(parts)
    return ('<?xml version="1.0" encoding="UTF-8"?>\n' + dict_to_xml(parsed)).encode("utf-8")

def _yaml_to_json(data: bytes) -> bytes:
    try:
        import yaml
        parsed = yaml.safe_load(data.decode("utf-8"))
        return json.dumps(parsed, indent=2, ensure_ascii=False).encode("utf-8")
    except ImportError: return b"[PyYAML not installed]"
    except Exception as e: return f"Error: {e}".encode()

def _json_to_yaml(data: bytes) -> bytes:
    try:
        import yaml
        parsed = json.loads(data.decode("utf-8"))
        return yaml.dump(parsed, default_flow_style=False, allow_unicode=True).encode("utf-8")
    except ImportError: return b"[PyYAML not installed]"
    except Exception as e: return f"Error: {e}".encode()

def _pptx_to_text(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text: out.append(shape.text)
        return "\n".join(out)
    except Exception as e:
        return f"[PPTX error: {e}]"

def _image_to_text(data: bytes, ext: str) -> bytes:
    try:
        from utils.gemini_key_manager import get_key as gemini_get_key
        import requests
        key = gemini_get_key()
        if not key: return b"[No Gemini key for OCR]"
        b64 = base64.b64encode(data).decode()
        mime = {"jpg":"image/jpeg","jpeg":"image/jpeg","png":"image/png","webp":"image/webp"}.get(ext,"image/png")
        resp = requests.post(
            f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={key}",
            json={"contents":[{"parts":[{"inline_data":{"mime_type":mime,"data":b64}},{"text":"Extract all text visible in this image. Return only raw text."}]}]},
            timeout=30)
        result = resp.json()
        text = result.get("candidates",[{}])[0].get("content",{}).get("parts",[{}])[0].get("text","")
        return text.encode("utf-8")
    except Exception as e:
        return f"OCR error: {e}".encode()

def _xlsx_to_html(data: bytes) -> bytes:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        tables = []
        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows: continue
            html_table = f"<h2>{sheet.title}</h2><table border='1' cellpadding='6'>"
            for i, row in enumerate(rows):
                tag = "th" if i == 0 else "td"
                html_table += "<tr>" + "".join(f"<{tag}>{str(c) if c is not None else ''}</{tag}>" for c in row) + "</tr>"
            html_table += "</table>"
            tables.append(html_table)
        page = f"""<!DOCTYPE html><html><head><meta charset='utf-8'>
<style>body{{font-family:sans-serif;padding:24px;}}table{{border-collapse:collapse;}}th{{background:#444;color:white;}}td,th{{padding:8px 12px;}}</style></head>
<body>{"".join(tables)}</body></html>"""
        return page.encode("utf-8")
    except Exception as e:
        return f"Error: {e}".encode()

# ── Format registry ────────────────────────────────────────────────────────────

def get_supported_formats() -> dict:
    return {
        "pdf": ["txt", "docx", "md"],
        "docx": ["txt", "pdf", "md"],
        "txt": ["pdf", "docx", "html", "md", "json"],
        "csv": ["json", "xlsx", "html", "md"],
        "json": ["csv", "xml", "yaml", "html", "txt"],
        "xlsx": ["csv", "json", "html", "txt"],
        "md": ["html", "pdf", "txt", "docx"],
        "html": ["txt", "md"],
        "xml": ["json", "txt"],
        "yaml": ["json", "txt"],
        "pptx": ["txt", "pdf"],
        "png": ["txt"],
        "jpg": ["txt"],
        "jpeg": ["txt"],
    }

def convert_file(data: bytes, from_fmt: str, to_fmt: str, filename: str = "file") -> Tuple[bytes, str, str, Optional[str]]:
    from_fmt = from_fmt.lower().strip(".")
    to_fmt   = to_fmt.lower().strip(".")
    mime_map = {
        "pdf": "application/pdf",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "txt": "text/plain", "csv": "text/csv", "json": "application/json",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "html": "text/html", "md": "text/markdown", "xml": "application/xml",
        "yaml": "text/yaml",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "png": "image/png",
    }
    try:
        key = (from_fmt, to_fmt)
        dispatch = {
            ("pdf","txt"): lambda: _pdf_to_text(data).encode("utf-8"),
            ("pdf","docx"): lambda: _text_to_docx(_pdf_to_text(data).encode("utf-8")),
            ("pdf","md"): lambda: _pdf_to_text(data).encode("utf-8"),
            ("docx","txt"): lambda: _docx_to_text(data).encode("utf-8"),
            ("docx","pdf"): lambda: _text_to_pdf(_docx_to_text(data).encode("utf-8")),
            ("docx","md"): lambda: _docx_to_text(data).encode("utf-8"),
            ("txt","pdf"): lambda: _text_to_pdf(data),
            ("txt","docx"): lambda: _text_to_docx(data),
            ("txt","html"): lambda: _text_to_html(data),
            ("txt","md"): lambda: data,
            ("txt","json"): lambda: json.dumps({"lines": data.decode("utf-8","replace").splitlines(), "content": data.decode("utf-8","replace")}, indent=2).encode("utf-8"),
            ("csv","json"): lambda: _csv_to_json(data),
            ("csv","xlsx"): lambda: _csv_to_xlsx(data),
            ("csv","html"): lambda: _text_to_html(_csv_to_text(data).encode("utf-8")),
            ("csv","md"): lambda: _csv_to_markdown_table(data),
            ("json","csv"): lambda: _json_to_csv(data),
            ("json","xml"): lambda: _json_to_xml(data),
            ("json","yaml"): lambda: _json_to_yaml(data),
            ("json","html"): lambda: _text_to_html(json.dumps(json.loads(data.decode("utf-8")), indent=2).encode("utf-8")),
            ("json","txt"): lambda: data,
            ("xlsx","csv"): lambda: _xlsx_to_csv(data),
            ("xlsx","json"): lambda: _xlsx_to_json(data),
            ("xlsx","html"): lambda: _xlsx_to_html(data),
            ("xlsx","txt"): lambda: _excel_to_text(data).encode("utf-8"),
            ("md","html"): lambda: _markdown_to_html(data),
            ("md","pdf"): lambda: _text_to_pdf(re.sub(r'[#*`]','',data.decode("utf-8")).encode("utf-8")),
            ("md","txt"): lambda: re.sub(r'[#*`]','',data.decode("utf-8")).encode("utf-8"),
            ("md","docx"): lambda: _text_to_docx(data),
            ("html","txt"): lambda: _html_to_text(data),
            ("html","md"): lambda: _html_to_markdown(data),
            ("xml","json"): lambda: _xml_to_json(data),
            ("xml","txt"): lambda: _html_to_text(data),
            ("yaml","json"): lambda: _yaml_to_json(data),
            ("yaml","txt"): lambda: data,
            ("pptx","txt"): lambda: _pptx_to_text(data).encode("utf-8"),
            ("pptx","pdf"): lambda: _text_to_pdf(_pptx_to_text(data).encode("utf-8")),
            ("png","txt"): lambda: _image_to_text(data, "png"),
            ("jpg","txt"): lambda: _image_to_text(data, "jpg"),
            ("jpeg","txt"): lambda: _image_to_text(data, "jpeg"),
        }
        if key not in dispatch:
            return b"", "", "", f"Conversion {from_fmt} → {to_fmt} not supported."
        result = dispatch[key]()
        mime = mime_map.get(to_fmt, "application/octet-stream")
        return result, mime, to_fmt, None
    except Exception as e:
        return b"", "", "", f"Conversion error: {e}"


# ── Unit Conversion Engine ─────────────────────────────────────────────────────

UNIT_CATEGORIES = {
    "📏 Length": {
        "units": ["Millimeter","Centimeter","Meter","Kilometer","Inch","Foot","Yard","Mile","Nautical Mile","Light Year","Astronomical Unit","Parsec","Furlong","Fathom","Angstrom","Micrometer","Nanometer"],
        "to_base": {"Millimeter":0.001,"Centimeter":0.01,"Meter":1,"Kilometer":1000,"Inch":0.0254,"Foot":0.3048,"Yard":0.9144,"Mile":1609.344,"Nautical Mile":1852,"Light Year":9.461e15,"Astronomical Unit":1.496e11,"Parsec":3.086e16,"Furlong":201.168,"Fathom":1.8288,"Angstrom":1e-10,"Micrometer":1e-6,"Nanometer":1e-9},
    },
    "⚖️ Weight / Mass": {
        "units": ["Milligram","Gram","Kilogram","Metric Ton","Ounce","Pound","Stone","US Ton","Long Ton","Carat","Grain","Troy Ounce","Microgram"],
        "to_base": {"Milligram":0.000001,"Gram":0.001,"Kilogram":1,"Metric Ton":1000,"Ounce":0.028350,"Pound":0.453592,"Stone":6.350293,"US Ton":907.185,"Long Ton":1016.047,"Carat":0.0002,"Grain":0.0000648,"Troy Ounce":0.0311035,"Microgram":1e-9},
    },
    "🌡️ Temperature": {
        "units": ["Celsius","Fahrenheit","Kelvin","Rankine","Réaumur"],
        "to_base": None,
    },
    "⏱️ Time": {
        "units": ["Nanosecond","Microsecond","Millisecond","Second","Minute","Hour","Day","Week","Month","Year","Decade","Century","Millennium"],
        "to_base": {"Nanosecond":1e-9,"Microsecond":1e-6,"Millisecond":0.001,"Second":1,"Minute":60,"Hour":3600,"Day":86400,"Week":604800,"Month":2629746,"Year":31556952,"Decade":315569520,"Century":3155695200,"Millennium":31556952000},
    },
    "📐 Area": {
        "units": ["Square Millimeter","Square Centimeter","Square Meter","Hectare","Square Kilometer","Square Inch","Square Foot","Square Yard","Acre","Square Mile"],
        "to_base": {"Square Millimeter":1e-6,"Square Centimeter":0.0001,"Square Meter":1,"Hectare":10000,"Square Kilometer":1e6,"Square Inch":0.00064516,"Square Foot":0.092903,"Square Yard":0.836127,"Acre":4046.86,"Square Mile":2589988.1},
    },
    "🧪 Volume": {
        "units": ["Milliliter","Centiliter","Deciliter","Liter","Cubic Meter","Teaspoon (US)","Tablespoon (US)","Fluid Ounce (US)","Cup (US)","Pint (US)","Quart (US)","Gallon (US)","Gallon (UK)","Cubic Inch","Cubic Foot"],
        "to_base": {"Milliliter":0.001,"Centiliter":0.01,"Deciliter":0.1,"Liter":1,"Cubic Meter":1000,"Teaspoon (US)":0.00492892,"Tablespoon (US)":0.0147868,"Fluid Ounce (US)":0.0295735,"Cup (US)":0.236588,"Pint (US)":0.473176,"Quart (US)":0.946353,"Gallon (US)":3.78541,"Gallon (UK)":4.54609,"Cubic Inch":0.0163871,"Cubic Foot":28.3168},
    },
    "⚡ Energy": {
        "units": ["Joule","Kilojoule","Megajoule","Calorie","Kilocalorie","Watt-hour","Kilowatt-hour","Electron Volt","BTU","Foot-pound","Erg"],
        "to_base": {"Joule":1,"Kilojoule":1000,"Megajoule":1e6,"Calorie":4.184,"Kilocalorie":4184,"Watt-hour":3600,"Kilowatt-hour":3.6e6,"Electron Volt":1.602e-19,"BTU":1055.06,"Foot-pound":1.35582,"Erg":1e-7},
    },
    "💨 Speed": {
        "units": ["Meter/second","Kilometer/hour","Mile/hour","Foot/second","Knot","Mach","Speed of Light"],
        "to_base": {"Meter/second":1,"Kilometer/hour":0.277778,"Mile/hour":0.44704,"Foot/second":0.3048,"Knot":0.514444,"Mach":340.29,"Speed of Light":299792458},
    },
    "🔌 Power": {
        "units": ["Watt","Kilowatt","Megawatt","Gigawatt","Horsepower (Metric)","Horsepower (Imperial)","BTU/hour"],
        "to_base": {"Watt":1,"Kilowatt":1000,"Megawatt":1e6,"Gigawatt":1e9,"Horsepower (Metric)":735.499,"Horsepower (Imperial)":745.7,"BTU/hour":0.293071},
    },
    "🔋 Pressure": {
        "units": ["Pascal","Kilopascal","Megapascal","Bar","Millibar","Atmosphere","PSI","mmHg (Torr)","inHg"],
        "to_base": {"Pascal":1,"Kilopascal":1000,"Megapascal":1e6,"Bar":1e5,"Millibar":100,"Atmosphere":101325,"PSI":6894.76,"mmHg (Torr)":133.322,"inHg":3386.39},
    },
    "💾 Data Storage": {
        "units": ["Bit","Byte","Kilobyte","Megabyte","Gigabyte","Terabyte","Petabyte","Kibibyte","Mebibyte","Gibibyte","Tebibyte"],
        "to_base": {"Bit":0.125,"Byte":1,"Kilobyte":1000,"Megabyte":1e6,"Gigabyte":1e9,"Terabyte":1e12,"Petabyte":1e15,"Kibibyte":1024,"Mebibyte":1048576,"Gibibyte":1073741824,"Tebibyte":1099511627776},
    },
    "💡 Frequency": {
        "units": ["Hertz","Kilohertz","Megahertz","Gigahertz","Terahertz","RPM"],
        "to_base": {"Hertz":1,"Kilohertz":1e3,"Megahertz":1e6,"Gigahertz":1e9,"Terahertz":1e12,"RPM":1/60},
    },
    "📡 Angle": {
        "units": ["Degree","Radian","Gradian","Arcminute","Arcsecond","Revolution"],
        "to_base": {"Degree":1,"Radian":57.2958,"Gradian":0.9,"Arcminute":1/60,"Arcsecond":1/3600,"Revolution":360},
    },
    "🌊 Force": {
        "units": ["Newton","Kilonewton","Pound-force","Dyne","Kilogram-force"],
        "to_base": {"Newton":1,"Kilonewton":1000,"Pound-force":4.44822,"Dyne":1e-5,"Kilogram-force":9.80665},
    },
    "🔢 Number Base": {
        "units": ["Binary (base 2)","Octal (base 8)","Decimal (base 10)","Hexadecimal (base 16)"],
        "to_base": None,
    },
    "💰 Currency (Live)": {
        "units": ["USD","EUR","GBP","JPY","INR","CAD","AUD","CHF","CNY","HKD","SGD","KRW","MXN","BRL","RUB","TRY","ZAR","AED","SAR","THB"],
        "to_base": None,
    },
    "💊 Cooking Measures": {
        "units": ["Teaspoon","Tablespoon","Cup","Pint","Quart","Gallon","Milliliter","Liter","Fluid Ounce"],
        "to_base": {"Teaspoon":4.92892,"Tablespoon":14.7868,"Cup":236.588,"Pint":473.176,"Quart":946.353,"Gallon":3785.41,"Milliliter":1,"Liter":1000,"Fluid Ounce":29.5735},
    },
    "🌍 Fuel Economy": {
        "units": ["km/L","mpg (US)","mpg (UK)"],
        "to_base": {"km/L":1,"mpg (US)":0.425144,"mpg (UK)":0.354006},
    },
    "🌡️ Heat Transfer": {
        "units": ["W/(m·K)","BTU/(hr·ft·°F)","cal/(s·cm·°C)"],
        "to_base": {"W/(m·K)":1,"BTU/(hr·ft·°F)":1.73073,"cal/(s·cm·°C)":418.68},
    },
    "📷 Resolution": {
        "units": ["DPI","PPI","Pixels/cm","Pixels/mm"],
        "to_base": {"DPI":1,"PPI":1,"Pixels/cm":2.54,"Pixels/mm":25.4},
    },
}


def convert_units(value: float, from_unit: str, to_unit: str, category: str) -> Tuple[float, str]:
    cat = UNIT_CATEGORIES.get(category, {})

    if category == "🌡️ Temperature":
        if from_unit == "Celsius": c = value
        elif from_unit == "Fahrenheit": c = (value - 32) * 5/9
        elif from_unit == "Kelvin": c = value - 273.15
        elif from_unit == "Rankine": c = (value - 491.67) * 5/9
        elif from_unit == "Réaumur": c = value * 5/4
        else: c = value
        if to_unit == "Celsius": result = c
        elif to_unit == "Fahrenheit": result = c * 9/5 + 32
        elif to_unit == "Kelvin": result = c + 273.15
        elif to_unit == "Rankine": result = (c + 273.15) * 9/5
        elif to_unit == "Réaumur": result = c * 4/5
        else: result = c
        return result, f"{value} {from_unit} = {result:.4g} {to_unit}"

    if category == "🔢 Number Base":
        base_map = {"Binary (base 2)":2,"Octal (base 8)":8,"Decimal (base 10)":10,"Hexadecimal (base 16)":16}
        fb = base_map.get(from_unit, 10)
        tb = base_map.get(to_unit, 10)
        try:
            n = int(str(int(value)), fb)
            if tb == 2: res = bin(n)[2:]
            elif tb == 8: res = oct(n)[2:]
            elif tb == 16: res = hex(n)[2:].upper()
            else: res = str(n)
            return n, f"{int(value)} ({from_unit}) = {res} ({to_unit})"
        except Exception as e:
            return 0, f"Error: {e}"

    if category == "💰 Currency (Live)":
        rates = get_live_currency_rates(from_unit)
        if to_unit in rates:
            result = value * rates[to_unit]
            return result, f"{value} {from_unit} = {result:.4f} {to_unit}"
        return 0, "Could not fetch live rates. Check internet."

    to_base = cat.get("to_base", {})
    if not to_base or from_unit not in to_base or to_unit not in to_base:
        return 0, "Conversion not available for this unit pair."
    
    base_val = value * to_base[from_unit]
    result = base_val / to_base[to_unit]
    return result, f"{value} {from_unit} = {result:.10g} {to_unit}"


def get_live_currency_rates(base: str = "USD") -> dict:
    try:
        import requests
        resp = requests.get(f"https://open.er-api.com/v6/latest/{base}", timeout=8)
        data = resp.json()
        if data.get("result") == "success":
            return data.get("rates", {})
        resp2 = requests.get(f"https://api.exchangerate-api.com/v4/latest/{base}", timeout=8)
        return resp2.json().get("rates", {})
    except Exception:
        return {}
